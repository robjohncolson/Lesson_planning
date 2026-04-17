"""Generate Day_2_Slides.pptx \u2014 classroom projection deck for Day 2.

One slide per phase with [Framework] + [DOK] badges, the prompt in large
font, and a time indicator. The slide displayed on student entry shows the
solo-paper Sign Flip task so late-arrivers see what to do immediately.

Designed for a traveling teacher whose AverTouch doccam is unreliable \u2014
every prompt students or evaluators might need is baked into the deck so
the teacher never has to write on a screen mid-class.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette (blue-family, matches pacer)
NAVY = RGBColor(0x0B, 0x3C, 0x6C)
BLUE = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
DARK = RGBColor(0x11, 0x22, 0x33)
GRAY = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)  # warm, used for transition cues


def add_background(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_badge(slide, text, left, top, *, color=NAVY, text_color=RGBColor(0xFF, 0xFF, 0xFF)):
    # Sized roughly to text length
    width = Inches(0.28 + 0.11 * len(text))
    height = Inches(0.35)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shp, width


def add_header(slide, phase_num, total, phase_title, framework, dok, minutes):
    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Phase title (left)
    add_text(slide, phase_title, Inches(0.5), Inches(0.18),
             Inches(9.5), Inches(0.9),
             size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Meta (below title)
    meta = f"Phase {phase_num}/{total}  \u00b7  Algebra 2  \u00b7  Lesson 3-5  \u00b7  Day 2"
    add_text(slide, meta, Inches(0.5), Inches(0.82),
             Inches(9), Inches(0.35),
             size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Badges (right)
    bx = Inches(10.3)
    by = Inches(0.4)
    if framework:
        shp, w = add_badge(slide, framework, bx, by, color=RGBColor(0x22, 0x33, 0x44))
        bx = bx + w + Inches(0.08)
    if dok and dok != "\u2014":
        shp, w = add_badge(slide, dok, bx, by, color=BLUE)

    # Minutes chip below badges
    if minutes is not None:
        mt = f"{minutes} min"
        add_text(slide, mt, Inches(10.3), Inches(0.85),
                 Inches(2.5), Inches(0.35),
                 size=14, bold=True,
                 color=RGBColor(0xCC, 0xDD, 0xEE),
                 align=PP_ALIGN.LEFT)


def add_phase_slide(prs, phase_num, total, phase_title, framework, dok,
                    minutes, body_lines, footer=None, body_size=24):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide)
    add_header(slide, phase_num, total, phase_title, framework, dok, minutes)

    # Body
    add_text(slide, body_lines,
             Inches(0.6), Inches(1.6),
             Inches(12.1), Inches(5.3),
             size=body_size, color=DARK)

    # Footer transition cue
    if footer:
        # Footer bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     0, Inches(6.85), SLIDE_W, Inches(0.65))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        add_text(slide, footer, Inches(0.5), Inches(6.92),
                 Inches(12.3), Inches(0.5),
                 size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    return slide


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---------- Slide 1: Title ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "ALGEBRA 2  \u00b7  LESSON 3-5",
             Inches(0.6), Inches(1.3), Inches(12), Inches(0.7),
             size=24, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, "Day 2  \u2014  Graphing from Factored Form",
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.4),
             size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, [
        "Math Objective: Identify zeros by factoring or synthetic division; use zeros to graph.",
        "Essential Question: How do the zeros and a point on the graph determine the equation?",
        "",
        "Materials: laptops (Desmos), blank paper, pencils, Blooket code.",
    ],
        Inches(0.6), Inches(3.8), Inches(12), Inches(3),
        size=20, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_text(s, "60-minute block  \u00b7  post-spring-break Monday",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             size=16, color=RGBColor(0x88, 0xAA, 0xCC))

    total = 8

    # ---------- Slide 2: Do Now A (solo paper) ----------
    add_phase_slide(prs, 1, total,
        "Do Now  \u2014  Sign Flip Prediction (Solo Paper)",
        "Do Now A", "DOK 2", 5, [
            "On the Do Now sheet in front of you:",
            "",
            "Compare   f(x) = x(x \u2212 4)(x + 3)",
            "  vs.      g(x) = \u2212x(x \u2212 4)(x + 3)",
            "",
            "1.  What will stay the EXACT SAME?",
            "2.  What will CHANGE? How do you know BEFORE graphing?",
            "3.  Finish: \u201cThe zeros stay the same because ___, and the graph changes because ___.\u201d",
            "",
            "Silent \u00b7 Pencil only \u00b7 No Desmos \u00b7 No packet yet",
        ],
        footer="When everyone's back:  Blooket code coming on the screen.")

    # ---------- Slide 3: Blooket Code ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, 2, total, "Blooket  \u2014  Log In", "Do Now B", "\u2014", 2)
    add_text(s, "Blooket code:",
             Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
             size=28, color=GRAY)
    # Big empty box for teacher to write/display code
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.5), Inches(2.6), Inches(10.3), Inches(3))
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = NAVY
    box.line.width = Pt(3)
    add_text(s, "( write the code here or paste a screenshot )",
             Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.5),
             size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "Chromebooks open \u00b7 Type the code \u00b7 Enter a nickname",
             Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(6.85), SLIDE_W, Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(s, "2 minutes to log in. Game waits for no one.",
             Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.5),
             size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # ---------- Slide 4: Blooket Game ----------
    add_phase_slide(prs, 3, total,
        "Blooket  \u2014  Warm-Up Game",
        "Do Now C", "DOK 1", 7, [
            "Topics:",
            "  \u2022  Finding zeros from factored form",
            "  \u2022  Counting intervals from zeros",
            "  \u2022  Difference of squares",
            "  \u2022  GCF + quadratic factoring",
            "  \u2022  Sign-flip primers",
            "",
            "Focus. Play fast. Teacher is watching the dashboard \u2014 which skills need a quick reset?",
        ],
        footer="Close game when timer hits 0.  \u201cPull out your Do Now sheet. Packet too.\u201d")

    # ---------- Slide 5: Sign Flip Synthesis ----------
    add_phase_slide(prs, 4, total,
        "Launch  \u2014  Sign Flip Synthesis",
        "Launch", "DOK 2", 6, [
            "Out: Do Now sheet  +  Student packet.",
            "",
            "1.  Turn and Talk with your partner:",
            "      What did you predict would stay the same?",
            "      What did you predict would change?",
            "",
            "2.  Class share  +  Desmos check.",
            "",
            "3.  Sentence frame:",
            "      \u201cThe zeros stay the same because ___.",
            "       The graph changes because ___.\u201d",
        ],
        footer="\u201cPens down. Practice #13 next. Factor, then graph in Desmos. Watch x = 4.\u201d")

    # ---------- Slide 6: Practice #13 ----------
    add_phase_slide(prs, 5, total,
        "Explore  \u2014  Practice #13 Discovery",
        "Explore", "DOK 2\u20133", 15, [
            "g(x) = x\u00b3 \u2212 8x\u00b2 + 16x",
            "",
            "Step 1.   Factor on blank paper.  (GCF first!)",
            "Step 2.   List the zeros.   How many DISTINCT zeros?",
            "Step 3.   Graph in Desmos.  Watch what happens at x = 4.",
            "Step 4.   Sketch on blank paper.",
            "",
            "In one sentence \u2014 what does the graph DO at x = 4?",
        ],
        footer="Capture student language (\u201ctouches,\u201d \u201cbounces\u201d). DO NOT name multiplicity today.")

    # ---------- Slide 7: Reverse Engineering ----------
    add_phase_slide(prs, 6, total,
        "Explore  \u2014  Build the Equation (Reverse)",
        "Explore", "DOK 3", 15, [
            "Write the equation of a polynomial with",
            "    zeros at  x = \u22122,  x = 1,  x = 4",
            "    passing through the point  (0, \u22128).",
            "",
            "Step 1.   Zeros \u2192 factors.   f(x) = a(    )(    )(    )",
            "Step 2.   Plug in (0, \u22128) to find a.",
            "Step 3.   Write the final equation.",
            "Step 4.   Justify to your partner in CER format (reference at top of packet).",
        ],
        footer="Push \u201cbecause\u201d / \u201csince.\u201d  If time short: verbal CER now, written CER as HW.")

    # ---------- Slide 8: Share / Summary ----------
    add_phase_slide(prs, 7, total,
        "Share / Summary",
        "Share/Summary", "DOK 2", 5, [
            "Essential Question:",
            "How do the zeros \u2014 and a single point on the graph \u2014",
            "determine the polynomial\u2019s equation?",
            "",
            "Language Objective check:",
            "Who used  \u201cThe zeros stay the same because...\u201d today?",
            "",
            "Self-rating on your packet:   \u2713  /  partly  /  not yet",
            "",
            "Preview Day 3:  tomorrow we name what happened at x = 4.",
        ],
        footer="\u201cPencils up. One question. CER format. 5 minutes.\u201d")

    # ---------- Slide 9: Exit Ticket ----------
    add_phase_slide(prs, 8, total,
        "Exit Ticket  \u2014  CER",
        "Exit Ticket", "DOK 2", 5, [
            "Savvas Practice #6:",
            "",
            "If you use zeros to sketch the graph of a polynomial function,",
            "how can you verify that your graph is correct?",
            "",
            "Write:   CLAIM  \u2192  EVIDENCE  \u2192  REASONING",
            "",
            "Use the CER reference at the top of your packet if you need it.",
            "Put your packet on the front desk as you leave.",
        ],
        footer="Packets to the front. Pack up. See you next class.",
        body_size=26)

    # ---------- Slide 10: Wrap-up ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "Nice work today.",
             Inches(0.6), Inches(2.3), Inches(12), Inches(1.2),
             size=52, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(s, [
        "Collect:  Exit Ticket  \u00b7  Do Now sheet  \u00b7  Reverse Engineering page",
        "",
        "Day 3 preview:  we name what happened at x = 4 in g(x).",
    ],
        Inches(0.6), Inches(4.1), Inches(12), Inches(2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF),
        align=PP_ALIGN.CENTER)

    prs.save(path)


if __name__ == "__main__":
    build("Day_2_Slides.pptx")
    print("Built Day_2_Slides.pptx")
