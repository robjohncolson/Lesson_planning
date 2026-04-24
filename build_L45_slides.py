"""L45_P1_Slides.pptx — projection deck for Lesson 4-5 (single-period quick).

Solving Rational Equations + DOK-3 Chemistry Alcohol Mixture capstone
(Topic 4 LEHS Q#2 + Q#6 + Q#7 prep). Compressed 2-example Launch
(Ex 1 solve + Ex 3 extraneous).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


# ── ID manifest (mirrors build_L45_P1_packets.py) ──────────────────────────

DO_NOW_ID   = "4-5-savvas-model-discuss-lesson-4-5-launch"
LAUNCH_IDS  = [
    "4-5-savvas-example-1-lesson-4-5",
    "4-5-savvas-example-3-lesson-4-5",
]
EXPLORE_IDS = [
    "4-5-savvas-try-it-1-lesson-4-5",
    "4-5-savvas-try-it-3-lesson-4-5",
    "4-5-savvas-q10",
    "4-5-savvas-q14",
    "4-5-savvas-q25",
    "4-5-savvas-q33",
]
REINFORCE_IDS = ["4-5-savvas-q8"]

_ALL_IDS = [DO_NOW_ID] + LAUNCH_IDS + EXPLORE_IDS + REINFORCE_IDS
qb.get_for_packet(_ALL_IDS)


# ── Drawing primitives ─────────────────────────────────────────────────────

def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, "Algebra 2  ·  Lesson 4-5  ·  Single Period (Quick)",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.8),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.85), Inches(11.7), Inches(0.8),
         size=24, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Student-centered · Topic 4 LEHS Q#2/Q#6/Q#7 rehearsal",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


RULES_CARD_LINES = [
    "1. Find the LCD of all fractions in the equation.",
    "2. Multiply BOTH SIDES by the LCD to clear fractions.",
    "3. Solve the resulting polynomial equation for candidates.",
    "4. CHECK each candidate in the ORIGINAL equation — any that makes a denominator zero is EXTRANEOUS.",
]


# ── P1 deck ────────────────────────────────────────────────────────────────

def build_p1(path="L45_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(
        prs,
        "Lesson 4-5 · Quick — Solving Rational Equations",
        "Solve + Extraneous Check + ⭐ Chemistry Mixture DOK-3 (Q#2/Q#6/Q#7 Prep)",
    )

    # ── Do Now ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Do Now — Model & Discuss: Conflicting Candidates",
           framework_tag="exploration", dok=1, minutes=5)
    card(s, f"{DO_NOW_ID}",
         ["Two students solved a rational equation and got DIFFERENT candidate solutions.",
          "",
          "A) If a candidate makes the ORIGINAL denominator zero, is it a real solution?",
          "B) Why might clearing denominators produce a value that doesn't work?",
          "C) Bridge: yesterday we ADDED rational expressions. Today we SOLVE rational EQUATIONS — what's the new trap?"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    # ── Launch ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Launch — Examples 1 & 3 + Rules [COMPRESSED]",
           framework_tag="teacher models", dok=2, minutes=13)
    card(s, "\U0001f4d8 Solving Rational Equations — 4-Step Procedure",
         RULES_CARD_LINES,
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), color=GREEN)
    card(s, "Examples: 4-5-savvas-example-1-lesson-4-5  &  4-5-savvas-example-3-lesson-4-5",
         ["Ex 1: SOLVE a rational equation — find LCD, clear, solve polynomial, check in ORIGINAL.",
          "Ex 3: EXTRANEOUS solution — state domain restrictions FIRST, solve, test each candidate.",
          "KEY: candidates that violate the original domain are extraneous, even if they satisfy the cleared equation."],
         Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.6), color=BLUE)

    # ── Explore ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Explore — Think-Pair-Share + ⭐ DOK-3 Chemistry",
           framework_tag="student work", dok="2-3", minutes=32)
    text(s, "6 items in order. Plan ~10 min for ⭐ Practice #33 (Chemistry Mixture DOK-3).",
         Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.1)
    items = [
        ("Try It 1 · 4-5-savvas-try-it-1-lesson-4-5",
         "Solve a rational equation. Find the LCD, clear, solve.", BLUE),
        ("Try It 3 · 4-5-savvas-try-it-3-lesson-4-5",
         "Identify the extraneous solution. Check in the ORIGINAL.", BLUE),
        ("Practice #10 · 4-5-savvas-q10",
         "Solve + check. Confirm candidates in original equation.", BLUE),
        ("Practice #14 · 4-5-savvas-q14",
         "Extraneous trap — state domain restrictions BEFORE clearing.", BLUE),
        ("Practice #25 · 4-5-savvas-q25",
         "Work-rate application: 1/a + 1/b = 1/t.", BLUE),
        ("Practice #33 ⭐ DOK-3 CHEMISTRY MIXTURE · 4-5-savvas-q33",
         "50 gal of 2% alcohol + x gal of 6% alcohol → 5% target. Solve for x. Part B: calculator TABLE verification. Topic 4 LEHS Q#6+Q#7 rehearsal. Plan ~10 min.", ACCENT),
    ]
    row_h = Inches(0.72)
    for lbl, txt, clr in items:
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), row_h, color=clr)
        top += row_h + Inches(0.05)

    # ── DOK-3 Spotlight ───────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "⭐ DOK-3 SPOTLIGHT — Chemistry Mixture (Practice #33)",
           framework_tag="DOK-3 driver", dok=3, minutes=10)
    card(s, "\U0001f9ea MIXTURE SETUP",
         ["50 gal of 2% alcohol solution. Add x gal of 6% alcohol.",
          "Resulting concentration:",
          "",
          "    f(x) = (50·0.02 + 0.06x) / (50 + x)",
          "",
          "PART A: Solve f(x) = 0.05 for x.",
          "",
          "PART B: Explain how to verify using a graphing calculator TABLE."],
         Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.7), color=ACCENT)
    card(s, "\U0001f4ac TEACHER NOTES",
         ["Move Part A wants: rational-equation setup → clear denominator → solve linear.",
          "",
          "Part B: enter y₁ = numerator expression, y₂ = 0.05(50+x). Use TABLE to find intersection.",
          "",
          "STUCK? Ask: \"What would x mean in gallons?\" BEFORE they compute — anchors reasonableness.",
          "",
          "(Answer hidden on student view — reveal only at Share.)",
          "",
          "BRIDGE: same move as Topic 4 LEHS Q#6 (solve); domain reasoning = Q#2; rate reasoning = Q#7."],
         Inches(6.7), Inches(1.55), Inches(6.1), Inches(5.7), color=GOLD)

    # ── Share / Summary ───────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Share / Summary — Chemistry Reveal + Q#2/Q#6/Q#7 Bridge",
           framework_tag="academic conversation", dok=2, minutes=5)
    card(s, "\U0001f4e2 Answer Key — Practice #33 (DOK-3 Chemistry)",
         ["(50·0.02 + 0.06x)/(50+x) = 0.05",
          "    → 1 + 0.06x = 0.05(50 + x)",
          "    → 1 + 0.06x = 2.5 + 0.05x",
          "    → 0.01x = 1.5",
          "    → x = 150 gallons",
          "",
          "Part B: y₁ = 50·0.02 + 0.06x;  y₂ = 0.05(50+x). TABLE shows y₁ = y₂ at x = 150.",
          "",
          "Sentence frame: \"I set up f(x) = 0.05, cleared the denominator, and solved for x to get ___ gallons. I verified with TABLE by entering y₁ = ___ and y₂ = ___.\"",
          "",
          "\U0001f517 Topic 4 LEHS bridge: Q#6 (solve), Q#2 (domain), Q#7 (rate) — all rehearsed."],
         Inches(0.6), Inches(1.55), Inches(12.1), Inches(5.7), color=BLUE)

    # ── Exit Ticket ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Exit Ticket — Summary Recap",
           framework_tag="not graded", dok="1-2", minutes=3)
    card(s, "✍️ Summary Exit — 3 stems",
         ["1. After clearing denominators in a rational equation, I must check each candidate solution in the ___ equation to catch ___ solutions.",
          "",
          "2. In Practice #33 (chemistry), the equation came from setting ___ equal to the target concentration 0.05.",
          "",
          "3. One biggest thing I learned today: ___________________________"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    prs.save(path)
    print(f"Built {path}")
    return path


if __name__ == "__main__":
    build_p1("L45_P1_Slides.pptx")
