"""L55_P1_Slides.pptx — projection deck for Lesson 5-5 (single-period standalone).

Function Operations + Composition Order DOK-3 (NOT on Topic 5 Assessment).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, "Algebra 2  ·  Lesson 5-5  ·  Single Period",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.8),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.85), Inches(11.7), Inches(0.8),
         size=24, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Klimsara-adapted · Student-centered · NOT on Topic 5 Assessment",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


RULES_CARD_LINES = [
    "(f±g)(x) = f(x)±g(x)",
    "(fg)(x) = f(x)·g(x)",
    "(f/g)(x) = f(x)/g(x),  g ≠ 0",
    "(f∘g)(x) = f(g(x))  —  g FIRST",
    "f∘g ≠ g∘f  in general",
]


# ── P1 deck ────────────────────────────────────────────────────────────────

def build_p1(path="L55_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(
        prs,
        "Lesson 5-5 · Quick — Function Operations + ⭐ Composition Order DOK-3 (Standalone)",
        "Function Operations: Add, Subtract, Multiply, Divide, Compose",
    )

    # ── Do Now ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Do Now — Profit = Revenue − Cost",
           framework_tag="exploration", dok=3, minutes=5)
    card(s, "5-5-savvas-model-discuss-lesson-5-5-launch",
         ["A company’s revenue R(x) and cost C(x) are both functions of items sold x.",
          "",
          "A) Write an expression for Profit P(x) using function notation.",
          "B) If R(x) = 3x + 10 and C(x) = x + 4, find P(x). Simplify.",
          "C) Bridge: If we can subtract functions, what else can we do with them?"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    # ── Launch ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Launch — Examples 1 & 4 + Rules",
           framework_tag="teacher models", dok=2, minutes=12)
    card(s, "\U0001f4d8 Function Operations — Rules",
         RULES_CARD_LINES,
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.4), color=GREEN)
    card(s, "Examples: 5-5-savvas-example-1-lesson-5-5  &  5-5-savvas-example-4-lesson-5-5",
         ["Ex 1: Add/Subtract/Multiply/Divide two functions. State domain restriction for division.",
          "Ex 4: Compose f∘g AND g∘f with concrete numbers first, then algebraically.",
          "KEY: Emphasize composition order — evaluate INNER function first."],
         Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.7), color=BLUE)

    # ── Explore ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Explore — Think-Pair-Share",
           framework_tag="student work", dok="2→3", minutes=33)
    text(s, "9 items in order. Plan ∼15 min for ⭐ Practice #32 (DOK-3 Standalone).",
         Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.1)
    items = [
        ("Try It 1 · 5-5-savvas-try-it-1-lesson-5-5",
         "Add and subtract two functions. Simplify the result.", BLUE),
        ("Try It 4 · 5-5-savvas-try-it-4-lesson-5-5",
         "Compose f∘g and g∘f. Are they equal?", BLUE),
        ("Practice #20 · 5-5-savvas-q20",
         "Find (f+g)(x). State the domain.", BLUE),
        ("Practice #21 · 5-5-savvas-q21",
         "Find (f−g)(x) and (fg)(x).", BLUE),
        ("Practice #24 · 5-5-savvas-q24",
         "Find (f/g)(x). State domain restriction.", BLUE),
        ("Practice #28 · 5-5-savvas-q28",
         "Evaluate a composition at a specific value.", BLUE),
        ("Practice #34 · 5-5-savvas-q34",
         "Mixed operations — apply multiple function operations.", BLUE),
        ("Practice #32 ⭐ DOK-3 STANDALONE · 5-5-savvas-q32",
         "Music store: 15%-off coupon vs $5-off. Which order saves more? JUSTIFY with BOTH compositions written out.", ACCENT),
        ("Reinforce · 5-5-savvas-q35",
         "SAT/ACT: evaluate f(g(5)). Assessment-adjacent — show all steps.", BLUE),
    ]
    row_h = Inches(0.62)
    for lbl, txt, clr in items:
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), row_h, color=clr)
        top += row_h + Inches(0.04)

    # ── Share / Summary ───────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Share / Summary — Composition Order Reveal",
           framework_tag="academic conversation", dok=2, minutes=5)
    card(s, "✅ Answer Key — Practice #32 (DOK-3)",
         ["(a) %-off first: 0.85(x − 5) = 0.85x − 4.25  ⇒  price = $72.25",
          "(b) $5-off first: 0.85x − 5  ⇒  at x=$90: $71.50",
          "(c) 15% FIRST saves MORE ($72.25 vs $71.50).",
          "    Reason: applying 15% to the larger pre-discount base saves more than applying it after the flat $5.",
          "",
          "Sentence frame: “The order matters because ___ — applying ___ first gives a lower price because ___."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=BLUE)

    # ── Exit Ticket ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Exit Ticket — f(g(5)) · SAT/ACT Reinforce",
           framework_tag="not graded", dok="1-2", minutes=3)
    card(s, "✏️ Practice #35 · 5-5-savvas-q35",
         ["Given f and g defined on your packet, evaluate f(g(5)).",
          "",
          "Steps:",
          "  1. Evaluate g(5) first.",
          "  2. Use that result as input for f.",
          "  3. Show both steps in writing.",
          "",
          "Note: Composition is NOT on the Topic 5 assessment — this is for SAT/ACT readiness."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.5), color=GOLD)

    prs.save(path)
    print(f"Built {path}")
    return path


if __name__ == "__main__":
    build_p1("L55_P1_Slides.pptx")
