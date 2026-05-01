"""L65_P1_Slides.pptx — projection deck for Lesson 6-5 (single-period).

Single deck: Properties of Logarithms + DOK-3 Paired Capstone (Power Property Proof + Richter Generalization).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, period_label, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, f"Algebra 2  ·  Lesson 6-5  ·  {period_label}",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
         size=28, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Student-centered",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


RULES_TEXT = (
    "PRODUCT: log_b(MN) = log_b M + log_b N.  "
    "QUOTIENT: log_b(M/N) = log_b M − log_b N.  "
    "POWER: log_b(M^n) = n·log_b M.  "
    "IDENTITY: log_b b = 1, log_b 1 = 0.  "
    "CHANGE OF BASE: log_b x = (log x)/(log b)."
)


# ── P1 deck (single period) ───────────────────────────────────────────────────

def build_p1(path="L65_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(
        prs,
        "Lesson 6-5 · Quick — Properties of Logarithms + ⭐ DOK-3 Proof & Richter",
        "Power Property Proof + Richter Generalization Paired Capstone",
    )

    # ── Do Now ────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Do Now — Error Analysis + Expand + Condense",
           framework_tag="bridge to launch", dok="1-2", minutes=7)
    card(s, "3 items — no calculator needed",
         ["#12: Change-of-base error analysis — find the mistake, name the rule violated.",
          "",
          "#4: Expand log₆(49/5) using Quotient Property. Write as difference of two logs.",
          "",
          "#5: Condense 5 ln s + 6 ln t into a single logarithm."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    # ── Launch ────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Launch — Example 3: Write as a Single Logarithm",
           framework_tag="teacher models", dok=2, minutes=5)
    card(s, "\U0001f4d8 Properties of Logarithms — Rules",
         [RULES_TEXT],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.0), color=GREEN)
    card(s, "Example: 6-5-savvas-example-3-lesson-6-5  — condensing activates all 3 properties",
         ["Order: Power Property first (coefficient → exponent), then Product or Quotient to merge.",
          "Condensing = working all three properties simultaneously.",
          "DOK-3 task (#13) will ask you to PROVE why Power Property works — from exponent laws.",
          "⚠️ Form B Q8 + Q15 (geometric series) are NOT on Topic 6 assessment."],
         Inches(0.6), Inches(3.85), Inches(12.1), Inches(3.3), color=BLUE)

    # ── Explore ───────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Explore — Think-Pair-Share + ⭐ DOK-3 Paired Capstone",
           framework_tag="student work", dok="2-3", minutes=28)
    text(s, "6 items (5 on Wed-F: skip #3). Plan ~12 min for ⭐ #13 + #40 DOK-3 Paired Capstone.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Practice #19 · 6-5-savvas-q19",
         "Condense: log₅ 6 + ½ log₅ y — Power Property first.", False),
        ("Practice #21 · 6-5-savvas-q21",
         "Condense: ⅓ ln 27 − 3 ln(2y) — simplify ⅓ ln 27 first.", False),
        ("Practice #34 · 6-5-savvas-q34",
         "Solve: 7^x = 100 using change of base.", False),
        ("Practice #3 · 6-5-savvas-q3  [SKIP on Wed-F 45-min]",
         "Amanda's expand error on log₄(c²d⁵) — identify the property she confused.", False),
        ("Practice #13  ⭐ DOK-3 PROOF · 6-5-savvas-q13",
         "Prove: log_b(M^n) = n·log_b M using exponent laws. "
         "Let x = log_b M, raise both sides to n, apply log definition.", True),
        ("Practice #40 Part C  ⭐ DOK-3 RICHTER · 6-5-savvas-q40",
         "One earthquake is 150× more intense than another. "
         "How much greater is its Richter magnitude? Generalize: R(150A₀) − R(A₀) = log 150 ≈ 2.18.", True),
    ]
    for lbl, txt, is_dok3 in items:
        color = ACCENT if is_dok3 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.78), color=color)
        top += Inches(0.84)

    # ── Share / Summary ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Share / Summary — DOK-3 Reveal",
           framework_tag="academic conversation", dok=2, minutes=5)
    card(s, "\U0001f4e2 Capstone Answers — #13 Proof + #40 Part C Richter",
         ["#13 PROOF: Let x = log_b M ⟹ b^x = M.",
          "Raise both sides to n: b^(nx) = M^n.",
          "By log definition: log_b(M^n) = nx = n·log_b M. ✓",
          "",
          "#40 Part C: R(150A₀) − R(A₀) = log(150A₀/S) − log(A₀/S) = log 150 ≈ 2.18.",
          "The larger earthquake is ~2.18 Richter units greater.",
          "",
          "Sentence frame: \"I proved the Power Property by ___, then used it to show that ___.\""],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=BLUE)

    # ── Exit Ticket ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Exit Ticket — Two-Part Summary",
           framework_tag="not graded", dok="1-2", minutes=5)
    card(s, "✍️ Exit Ticket",
         ["A biologist models a population as P = 3 log t + log 4.",
          "",
          "(a) Write P as a single logarithm: P = log(___)",
          "(b) Solve for t when P = 2: t = ___",
          "",
          "Expected: (a) P = log(4t³)    (b) t = (25)^(1/3) ≈ 2.924",
          "",
          "One biggest thing I learned today: ___."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


if __name__ == "__main__":
    p = build_p1("L65_P1_Slides.pptx")
    print(f"Built {p}")
