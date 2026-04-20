"""L41_P{1,2,3}_Slides.pptx — projection decks for Lesson 4-1 Klimsara-adapted.

Three decks, one per period. Same slide helpers as build_L35_slides.py.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, period_label, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, f"Algebra 2  ·  Lesson 4-1  ·  {period_label}",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
         size=28, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Klimsara-adapted · Student-centered · No Blooket",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


# ── P1 deck ───────────────────────────────────────────────────────────────

def build_p1(path="L41_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 4-1 · Period 1",
                "Inverse Variation — Introduction")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Do Now — Notice & Wonder",
           framework_tag="exploration", dok=2, minutes=5)
    card(s, "🔲 Two Rectangles · Both have area = 144 square units",
         ["Rectangle 1:  length = 72,  width = 2",
          "Rectangle 2:  length = 24,  width = 6",
          "",
          "A) Notice: 1 thing you notice",
          "B) Wonder: 1 thing you wonder",
          "C) Sketch: a third rectangle with area 144"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Launch — Inverse Variation",
           framework_tag="teacher models", dok=2, minutes=12)
    card(s, "📘 Inverse Variation Rules",
         ["• Two variables vary INVERSELY when their product is constant:",
          "        xy = k        or        y = k/x",
          "• k is the CONSTANT OF VARIATION",
          "• As one increases, the other decreases proportionally",
          "",
          "Example: x = 10, y = 3  →  k = xy = 30  →  y = 30/x",
          "When x = −6:  y = 30/(−6) = −5",
          "",
          "⚠️ COMMON ERROR: y/x = k is DIRECT variation (not inverse)"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GREEN)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Explore — Think-Pair-Share",
           framework_tag="student work", dok="1-2", minutes=33)
    text(s, "7 items in order. Use the Inverse Variation Rules.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.1)
    items = [
        ("Try It 1 · Identify inverse variation", "Check two tables — compute xy for each column."),
        ("Practice #9 · Generalize", "How can you tell from the table that it's NOT an inverse variation?"),
        ("Practice #14 · Another table check", "x: −1/4, −1/2, 1/3, 2, 5, 11 ; y: −9/2, −9, 6, 36, 90, 198."),
        ("Try It 2 · Write the equation", "x = 6 when y = 1/2. Find the equation and y at x = 15."),
        ("Practice #13 · Generalize", "Write k in terms of x and y."),
        ("Practice #16 · Apply", "x = 3 when y = 2/3. Find y when x = −1."),
        ("Practice #10 · Construct Arguments", "Explain why zero cannot be in the domain."),
    ]
    for lbl, txt in items:
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.65), color=BLUE)
        top += Inches(0.72)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Share / Summary",
           framework_tag="academic conversation", dok=2, minutes=5)
    card(s, "📢 Sentence frame",
         ["\"As ___ increases, ___ decreases.",
          " Their product is ___, so the equation is y = ___/x.\""],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.8), color=BLUE)
    card(s, "🔭 Preview of Period 2",
         ["Next class: solve real-world inverse variation problems.",
          "You'll distinguish direct from inverse in a single scenario — the Ramón Performance Task."],
         Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.3), color=ACCENT)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Exit Ticket — Summary Recap",
           framework_tag="not graded", dok="1-2", minutes=2)
    card(s, "✍️ Complete on your exit ticket",
         ["1. Today I learned that ___.",
          "2. The difference between inverse and direct variation is ___.",
          "3. One thing I'm still unsure about is ___."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


# ── P2 deck ───────────────────────────────────────────────────────────────

def build_p2(path="L41_P2_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 4-1 · Period 2",
                "Inverse Variation Applications + DOK-3 Performance Task")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Do Now — Is this inverse variation?",
           framework_tag="bridge from P1", dok=1, minutes=5)
    card(s, "Practice #15 · Table check",
         ["x: 1, 2, 3, 4, 5, 6",
          "y: 60, 30, 20, 15, 12, 10",
          "",
          "Is this an inverse variation? Compute xy for each column.",
          "If yes, state k."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Launch — Example 3 (Bouzouki)",
           framework_tag="teacher models", dok=2, minutes=12)
    card(s, "📘 Application Rules for Inverse Variation",
         ["1. Identify two varying quantities and their constant PRODUCT k.",
          "2. Compute k from one given pair.",
          "3. Write the equation y = k/x (rename variables to fit).",
          "4. Solve for the unknown.",
          "",
          "Example (bouzouki string): length s varies inversely with frequency f.",
          "  • 26-inch string → 329.63 cycles/sec → k = 8,570.38",
          "  • 13-inch string → f = 8,570.38 / 13 = 659.26 cycles/sec",
          "",
          "⚠️ DIRECT (y/x = k) vs INVERSE (xy = k) — both appear in today's DOK-3!"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GREEN)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Explore — TPS + DOK-3 Performance Task",
           framework_tag="student work", dok="2-3", minutes=33)
    text(s, "5 items in order. Plan ~15 min for Practice #26 (Ramón).",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Try It 3 · Ice cube melting", "Time to melt varies inversely with air temperature."),
        ("Practice #12 · HOT: inverse SQUARE", "y = k/x². If x is multiplied by 4, what happens to y?"),
        ("Practice #17 · Radio wavelength (graph)", "w varies inversely with frequency f. Find f when w = 375 m."),
        ("Practice #23 · Boyle's Law", "Volleyball 300 in³ at 4.5 psi vs basketball 415 in³ at 8 psi."),
        ("Practice #26  ⭐ PERFORMANCE TASK", "Ramón's road trip — Part A is DIRECT, Part B is INVERSE."),
    ]
    for i, (lbl, txt) in enumerate(items):
        color = ACCENT if i == 4 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.9), color=color)
        top += Inches(1.0)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Share / Summary",
           framework_tag="DOK-3 reveal + Wed bridge", dok=2, minutes=5)
    card(s, "📢 Sentence frame · Ramón",
         ["\"Part A is ___ variation because distance and time have a constant RATIO of ___.",
          " Part B is ___ variation because gas and time have a constant PRODUCT of ___.\""],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), color=BLUE)
    card(s, "🔭 Bridge to Period 3 (Wednesday)",
         ["Tomorrow we graph y = 1/x and its translations.",
          "Watch for: asymptotes, domain, range, intercepts.",
          "Period 3 is ASSESSMENT-CRITICAL for the Topic 4 8-Q assessment."],
         Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), color=ACCENT)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Exit Ticket — Summary Recap",
           framework_tag="not graded", dok="1-2", minutes=3)
    card(s, "✍️ Complete on your exit ticket",
         ["1. Today I learned that ___.",
          "2. The biggest trap in Ramón's task was ___.",
          "3. I distinguish direct from inverse variation by ___."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


# ── P3 deck ───────────────────────────────────────────────────────────────

def build_p3(path="L41_P3_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 4-1 · Period 3",
                "Reciprocal Function + Translations (Assessment-Critical)")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Do Now — Generalize + Predict",
           framework_tag="bridge from P2", dok=1, minutes=5)
    card(s, "Practice #13 · Generalize",
         ["Write k in terms of x and y for an inverse variation.",
          "",
          "Then 30-second prediction:",
          "  What does the graph of y = 1/x LOOK like?",
          "  (No wrong answers — just predict.)"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Launch — Example 4 + Example 5",
           framework_tag="teacher models 2 examples", dok=2, minutes=12)
    card(s, "📘 Parent Reciprocal Function · f(x) = 1/x",
         ["• Vertical asymptote: x = 0        • Horizontal asymptote: y = 0",
          "• Domain: {x | x ≠ 0}              • Range: {y | y ≠ 0}",
          "• Graph: hyperbola in Quadrants I and III"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.4), color=GREEN)
    card(s, "📘 Translated Reciprocal · g(x) = 1/(x − h) + k",
         ["• Vertical asymptote shifts to  x = h",
          "• Horizontal asymptote shifts to  y = k",
          "• Center (asymptote intersection): (h, k)",
          "• Example: g(x) = 1/(x − 3) + 2 → asymptotes x = 3, y = 2",
          "",
          "y-intercept: set x = 0 → y = −1/h + k        x-intercept: set y = 0 → x = h − 1/k"],
         Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.9), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Explore — TPS · Assessment Rehearsal",
           framework_tag="student work · pure DOK-2", dok=2, minutes=35)
    text(s, "6 items. Plan ~8 min for ⭐ Practice #19 (intercepts + asymptotes).",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Try It 4 · Graph g(x) = 10/x", "Parent scaled by 10. State asymptotes, domain, range."),
        ("Practice #11 · Error Analysis", "y = 5/x · Describe the labeling error for (2, 4) and (−2, −4)."),
        ("Practice #18 · Graph y = −2/x", "Reflected parent. State asymptotes, domain, range."),
        ("Try It 5 · Graph g(x) = 1/(x+2) − 4", "Translated. Identify (h, k), asymptotes, domain, range."),
        ("Practice #19  ⭐ ASSESSMENT-CRITICAL", "g(x) = 1/(x−2) + 6. Asymptotes + domain/range + INTERCEPTS."),
        ("Practice #21 · Video game storage", "160 games at 2.0 GB. Write equation, complete table, graph."),
    ]
    for i, (lbl, txt) in enumerate(items):
        color = ACCENT if i == 4 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.8), color=color)
        top += Inches(0.85)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Share / Summary · Concept Summary",
           framework_tag="LEHS 8-Q preview", dok=2, minutes=5)
    card(s, "📊 Concept Summary — Lesson 4-1",
         ["INVERSE VARIATION (y = k/x):",
          "  As one variable increases, the other decreases. k is constant.",
          "",
          "TRANSLATED RECIPROCAL (y = a/(x − h) + k):",
          "  Asymptotes shift to x = h and y = k.",
          "  (h, k) is the new center.",
          "",
          "🎯 Topic 4 8-Q assessment: expect Q#3 on asymptotes, Q#5 on translation descriptions."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=BLUE)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Exit Ticket — Pre-Assessment Confidence",
           framework_tag="confidence check", dok="1-2", minutes=3)
    card(s, "✍️ Rate yourself 1–5 + name one review item",
         ["• I can identify asymptotes of y = 1/(x−h) + k.  Confidence: 1 / 2 / 3 / 4 / 5",
          "• I can describe a translation from f(x) = 1/x to g(x).  Confidence: 1 / 2 / 3 / 4 / 5",
          "",
          "One thing I want to review before Topic 4 assessment: ___"],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


if __name__ == "__main__":
    for fn in [build_p1, build_p2, build_p3]:
        p = fn()
        print(f"Built {p}")
