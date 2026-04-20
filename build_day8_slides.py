"""Day_8_Slides.pptx — Day 8 projection deck (Quiz + Stations).

Six slides: Title + 5 phase slides (Do Now A / Quiz / Transition /
Stations / Exit).  No Blooket slide.  Lighter deck than Day 2-3.

Adapted from build_day23_slides.py.
"""

import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x3C, 0x6C)
BLUE = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
DARK = RGBColor(0x11, 0x22, 0x33)
GRAY = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
TEAL = RGBColor(0x00, 0x97, 0xA7)
GREEN = RGBColor(0x1B, 0x5E, 0x20)


def add_background(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_badge(slide, text, left, top, *, color=NAVY,
              text_color=RGBColor(0xFF, 0xFF, 0xFF)):
    width = Inches(0.28 + 0.11 * len(text))
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, Inches(0.35))
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shp, width


def add_header(slide, phase_num, total, phase_title, framework, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_text(slide, phase_title, Inches(0.5), Inches(0.18),
             Inches(9.5), Inches(0.9),
             size=34, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    meta = f"Phase {phase_num}/{total}  \u00b7  Algebra 2  \u00b7  Lesson 3-5  \u00b7  Day 8"
    add_text(slide, meta, Inches(0.5), Inches(0.82),
             Inches(9), Inches(0.35),
             size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    bx = Inches(10.3)
    by = Inches(0.4)
    if framework:
        shp, w = add_badge(slide, framework, bx, by,
                           color=RGBColor(0x22, 0x33, 0x44))
        bx = bx + w + Inches(0.08)
    if dok and dok != "\u2014":
        add_badge(slide, dok, bx, by, color=BLUE)

    if minutes is not None:
        add_text(slide, f"{minutes} min", Inches(10.3), Inches(0.85),
                 Inches(2.5), Inches(0.35),
                 size=14, bold=True,
                 color=RGBColor(0xCC, 0xDD, 0xEE))


def add_phase_slide(prs, phase_num, total, phase_title, framework, dok,
                    minutes, body_lines, footer=None, body_size=22):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, phase_num, total, phase_title, framework, dok, minutes)
    add_text(slide, body_lines,
             Inches(0.6), Inches(1.6),
             Inches(12.1), Inches(5.3),
             size=body_size, color=DARK)
    if footer:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     0, Inches(6.85), SLIDE_W, Inches(0.65))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        add_text(slide, footer, Inches(0.5), Inches(6.92),
                 Inches(12.3), Inches(0.5),
                 size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    return slide


def _fmt(text):
    SUPS = {"2": "\u00b2", "3": "\u00b3", "4": "\u2074", "5": "\u2075"}
    out, i = [], 0
    text = str(text)
    while i < len(text):
        ch = text[i]
        if ch == "^" and i + 1 < len(text) and text[i+1] in SUPS:
            out.append(SUPS[text[i+1]])
            i += 2
            continue
        out.append("\u2212" if ch == "-" else ch)
        i += 1
    return "".join(out)


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ----------------------------------------------------------------
    # Title slide
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "ALGEBRA 2  \u00b7  LESSON 3-5",
             Inches(0.6), Inches(1.3), Inches(12), Inches(0.7),
             size=22, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, "Day 8  \u2014  Lesson Quiz + Differentiated Stations",
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.4),
             size=38, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, [
        "Do Now A (vocabulary)  \u00b7  Quiz (25 min, silent)  \u00b7  Stations  \u00b7  Exit",
        "",
        "All quiz items from Savvas Lesson 3-5 Lesson Quiz.",
        "No Blooket today.",
    ],
        Inches(0.6), Inches(3.6), Inches(12), Inches(3.2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_text(s, "55-min lesson  \u00b7  Savvas-only",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             size=16, color=RGBColor(0x88, 0xAA, 0xCC))

    total = 5

    # ----------------------------------------------------------------
    # Phase 1 — Do Now A
    # ----------------------------------------------------------------
    do_now_q = qb.get("3-5-savvas-q28")
    prompt_snip = _fmt(do_now_q["prompt"])[:160] if do_now_q else "Vocabulary warm-up (q28)"
    add_phase_slide(prs, 1, total,
        "Do Now \u2014 Vocabulary Warm-Up",
        "Do Now A", "DOK 1", 5, [
            "On the Do Now sheet \u2014 Savvas Practice (bank: 3-5-savvas-q28):",
            "",
            f"  {prompt_snip}",
            "",
            "Complete both blanks.  Work alone.",
            "",
            "When done: re-read your notes on Lesson 3-5.",
        ],
        footer="5 minutes. Silent. Pencil only. No Desmos.")

    # ----------------------------------------------------------------
    # Phase 2 — Quiz
    # ----------------------------------------------------------------
    quiz_items = qb.get_for_packet(
        ["3-5-lq-q1a","3-5-lq-q1b","3-5-lq-q2","3-5-lq-q3","3-5-lq-q4","3-5-lq-q5"]
    )
    quiz_body = [
        "Savvas Lesson Quiz  (6 items, 25 minutes):",
        "",
        f"  Q1a.  {_fmt(quiz_items[0]['prompt'])[:90]}",
        f"  Q1b.  {_fmt(quiz_items[1]['prompt'])[:90]}",
        f"  Q2.   {_fmt(quiz_items[2]['prompt'])[:90]}",
        f"  Q3.   {_fmt(quiz_items[3]['prompt'])[:90]}",
        f"  Q4.   {_fmt(quiz_items[4]['prompt'])[:90]}",
        f"  Q5.   {_fmt(quiz_items[5]['prompt'])[:90]}",
    ]
    add_phase_slide(prs, 2, total,
        "Quiz \u2014 Savvas Lesson Quiz",
        "Assessment", "DOK 1\u20132", 25, quiz_body,
        footer="Silent. Pencil only. No Desmos. Show all work.",
        body_size=19)

    # ----------------------------------------------------------------
    # Phase 3 — Transition + Stations intro
    # ----------------------------------------------------------------
    add_phase_slide(prs, 3, total,
        "Stations \u2014 Differentiated Practice",
        "Stations", "DOK 2", 18, [
            "Rate yourself  1 / 2 / 3  on how the quiz felt.",
            "",
            "  \u25ba  1 (needed help)    \u2192  Support station  (re-teach: factoring + multiplicity)",
            "  \u25ba  2 (mostly there)  \u2192  Core station      (sign chart + complex zeros)",
            "  \u25ba  3 (confident)     \u2192  Extend station    (firework modeling, DOK 3)",
            "",
            "Circle your station(s) on the student packet.",
            "You have 18 minutes. You may do 2 stations.",
        ],
        footer="Student choice. Support \u2192 Core \u2192 Extend.")

    # ----------------------------------------------------------------
    # Phase 4 — Station detail (Support / Core / Extend)
    # ----------------------------------------------------------------
    # Resolve station IDs inline (same logic as build_day8_packets)
    _SUPPORT_ITEM_1 = "3-5-savvas-q14"
    _support_dyn = qb.select(lesson="3-5", dok=1, tags=["day3-multiplicity"], limit=2)
    _support_dyn = [q for q in _support_dyn if q["id"] != _SUPPORT_ITEM_1]
    _SUPPORT_ITEM_2 = _support_dyn[0]["id"] if _support_dyn else None
    _CORE_ITEM_1 = "3-5-savvas-q17"
    _CORE_ITEM_2 = "3-5-savvas-q22"
    _EXTEND_ITEM = "3-5-savvas-q25" if qb.get("3-5-savvas-q25") else None
    ext_label = _EXTEND_ITEM or "(not in bank \u2014 see teacher)"
    add_phase_slide(prs, 4, total,
        "Station Tasks",
        "Stations", "DOK 1\u20133", 18, [
            f"SUPPORT  [{_SUPPORT_ITEM_1}, {_SUPPORT_ITEM_2 or 'DOK 1'}]",
            "  Find zeros, state multiplicity, describe behavior (crosses / tangent).",
            "",
            f"CORE  [{_CORE_ITEM_1}  +  {_CORE_ITEM_2}]",
            "  Real + complex zeros; polynomial inequality sign chart.",
            "",
            f"EXTEND  [{ext_label}]  \u2014 DOK 3",
            "  Firework height model: domain, zeros in context, vertex.",
            "",
            "All items from Savvas bank. Show full work on the student packet.",
        ],
        footer="Teacher circulates. Questions only \u2014 no answers given.",
        body_size=20)

    # ----------------------------------------------------------------
    # Phase 5 — Exit
    # ----------------------------------------------------------------
    add_phase_slide(prs, 5, total,
        "Exit \u2014 Summary Recap",
        "Exit", "DOK 1\u20132", 5, [
            "Fill in the blanks on your packet (no notes, no Desmos):",
            "",
            "1.  A zero of a polynomial is an x-value where f(x) = ___.  "
            "The graph ___ or is ___ at that point.",
            "",
            "2.  EVEN multiplicity \u2192 graph is ___ to the x-axis (turning point).",
            "",
            "3.  To solve a polynomial inequality: find ___ and build a ___ chart.",
            "",
            "4.  My biggest takeaway from Lesson 3-5 (one sentence):",
            "    Place packets on the front desk as you leave.",
        ],
        footer="Packets to the front. Pack up.",
        body_size=21)

    # ----------------------------------------------------------------
    # Wrap slide
    # ----------------------------------------------------------------
    sw = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(sw, NAVY)
    add_text(sw, "Great work.",
             Inches(0.6), Inches(2.3), Inches(12), Inches(1.2),
             size=50, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(sw, [
        "Collect:  Do Now sheets  \u00b7  Quiz packets  \u00b7  Student packets",
        "",
        "Lesson 3-5 complete. On to Lesson 3-6.",
    ],
        Inches(0.6), Inches(4.1), Inches(12), Inches(2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF),
        align=PP_ALIGN.CENTER)

    prs.save(path)
    print(f"  Built {path}")


if __name__ == "__main__":
    build("Day_8_Slides.pptx")
    print()
    print("Built Day_8_Slides.pptx")
