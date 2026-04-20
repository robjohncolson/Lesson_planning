"""Day_23_Slides.pptx \u2014 Combined Day 2-3 projection deck (Savvas-only).

Seven phase slides + title + wrap. One slide per phase with framework +
DOK badges. All work items trace to Savvas bank IDs \u2014 the projection deck
shows the same items as the packet so students scanning the screen see
the same problems they're working.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x3C, 0x6C)
BLUE = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
DARK = RGBColor(0x11, 0x22, 0x33)
GRAY = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)


def add_background(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_badge(slide, text, left, top, *, color=NAVY,
              text_color=RGBColor(0xFF, 0xFF, 0xFF)):
    width = Inches(0.28 + 0.11 * len(text))
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, Inches(0.35))
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shp, width


def add_header(slide, phase_num, total, phase_title, framework, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_text(slide, phase_title, Inches(0.5), Inches(0.18),
             Inches(9.5), Inches(0.9),
             size=34, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    meta = f"Phase {phase_num}/{total}  \u00b7  Algebra 2  \u00b7  Lesson 3-5  \u00b7  Day 2-3"
    add_text(slide, meta, Inches(0.5), Inches(0.82),
             Inches(9), Inches(0.35),
             size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    bx = Inches(10.3)
    by = Inches(0.4)
    if framework:
        shp, w = add_badge(slide, framework, bx, by,
                           color=RGBColor(0x22, 0x33, 0x44))
        bx = bx + w + Inches(0.08)
    if dok and dok != "\u2014":
        add_badge(slide, dok, bx, by, color=BLUE)

    if minutes is not None:
        add_text(slide, f"{minutes} min", Inches(10.3), Inches(0.85),
                 Inches(2.5), Inches(0.35),
                 size=14, bold=True,
                 color=RGBColor(0xCC, 0xDD, 0xEE))


def add_phase_slide(prs, phase_num, total, phase_title, framework, dok,
                    minutes, body_lines, footer=None, body_size=22):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, phase_num, total, phase_title, framework, dok, minutes)
    add_text(slide, body_lines,
             Inches(0.6), Inches(1.6),
             Inches(12.1), Inches(5.3),
             size=body_size, color=DARK)
    if footer:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     0, Inches(6.85), SLIDE_W, Inches(0.65))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        add_text(slide, footer, Inches(0.5), Inches(6.92),
                 Inches(12.3), Inches(0.5),
                 size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    return slide


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "ALGEBRA 2  \u00b7  LESSON 3-5",
             Inches(0.6), Inches(1.3), Inches(12), Inches(0.7),
             size=22, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, "Combined Day 2-3  \u2014  Zeros + Multiplicity",
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.4),
             size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, [
        "Essential Question:",
        "How does the factored form of a polynomial tell you everything",
        "you need to sketch its graph \u2014 zeros AND behavior at each zero?",
        "",
        "Every work item today comes from Savvas.",
    ],
        Inches(0.6), Inches(3.6), Inches(12), Inches(3.2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_text(s, "55-min lesson  \u00b7  Savvas-only (no fabricated items)",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             size=16, color=RGBColor(0x88, 0xAA, 0xCC))

    total = 7

    add_phase_slide(prs, 1, total,
        "Do Now  \u2014  Zeros + Multiplicities",
        "Do Now A", "DOK 1", 5, [
            "On the Do Now sheet \u2014  SAVVAS LESSON QUIZ Q4:",
            "",
            "\u201cWhat are the zeros and their multiplicities for",
            "y = x\u00b3 + 3x\u00b2 + x + 3, shown in the graph?\u201d",
            "",
            "Read directly from the factored form of the graph.",
            "Use multiplicity language: even \u2192 touch, odd \u2192 cross.",
            "",
            "Silent  \u00b7  Pencil only  \u00b7  No Desmos",
        ],
        footer="When everyone\u2019s in:  Blooket code on the screen.")

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s2)
    add_header(s2, 2, total, "Blooket  \u2014  Log In", "Do Now B", "\u2014", 2)
    add_text(s2, "Blooket code:",
             Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
             size=28, color=GRAY)
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.5), Inches(2.6), Inches(10.3), Inches(3))
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = NAVY
    box.line.width = Pt(3)
    add_text(s2, "( write the code here )",
             Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.5),
             size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s2, "Chromebooks open  \u00b7  Type the code  \u00b7  Enter nickname",
             Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.85),
                              SLIDE_W, Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(s2, "2 minutes to log in. Game waits for no one.",
             Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.5),
             size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    add_phase_slide(prs, 3, total,
        "Blooket  \u2014  Rule Recall",
        "Do Now C", "DOK 1", 7, [
            "Pure rule recall.  The rules you\u2019ll need in 30 minutes:",
            "",
            "  \u2022  Zero Product Property  (zero \u21d4 factor = 0)",
            "  \u2022  Sign direction  (zero at x = c \u2192 factor (x \u2212 c))",
            "  \u2022  Multiplicity = exponent on the factor",
            "  \u2022  EVEN  \u2192  TOUCH and turn",
            "  \u2022  ODD   \u2192  CROSS",
            "  \u2022  Highest power  \u2192  end behavior",
            "",
            "Play fast.  Teacher is watching the dashboard.",
        ],
        footer="Close game at 0.  \u201cPull out your Do Now. Packet too.\u201d")

    add_phase_slide(prs, 4, total,
        "Launch  \u2014  Savvas Example 2",
        "Launch", "DOK 2", 12, [
            "Graph both in Desmos.  Watch  x = \u22122:",
            "",
            "    (1)  y = (x + 2)\u00b2 (x \u2212 1)",
            "    (2)  y = (x + 2)\u00b3 (x \u2212 1)",
            "",
            "THE WORD:  the number of times a factor appears is its MULTIPLICITY.",
            "",
            "Fill the T-chart on your packet:  ODD  vs.  EVEN.",
            "Sentence frame: \u201cThe multiplicity is ___, so the graph ___ at x = ___.\u201d",
        ],
        footer="\u201cPackets to Practice. Five Savvas items. A + B required, "
               "one of C/D/E.\u201d")

    add_phase_slide(prs, 5, total,
        "Practice  \u2014  Savvas Try Its + #14/15/16",
        "Practice", "DOK 2", 20, [
            "For each:  list zeros, multiplicity of each, and behavior.",
            "",
            "  A.  f(x) = x(x + 4)(x \u2212 1)\u2074        [Try It 2a]",
            "  B.  f(x) = (x\u00b2 + 9)(x \u2212 1)\u2075(x + 2)\u00b2    [Try It 2b]",
            "  C.  f(x) = x\u00b3 \u2212 8x\u00b2 + 16x           [Practice #14]",
            "  D.  g(x) = x\u00b3 \u2212 x\u00b2 \u2212 25x + 25         [Practice #15]",
            "  E.  f(x) = 9x\u2074 \u2212 40x\u00b2 + 16           [Practice #16]",
            "",
            "A + B required.  Pick one of C/D/E.  Fast finishers: do all five.",
        ],
        footer="Teacher circulates.  Prompts only \u2014 answer questions with questions.")

    add_phase_slide(prs, 6, total,
        "Share / Summary",
        "Share/Summary", "DOK 2", 4, [
            "Essential Question:",
            "How does the factored form tell you everything you need to sketch?",
            "",
            "Self-rating on your packet:   \u2713  /  partly  /  not yet",
            "  1. Read multiplicity from factored form",
            "  2. Predict cross vs. touch from multiplicity",
            "  3. Factor a polynomial in standard form",
            "",
            "Preview Day 4-5: real + complex zeros + Savvas storage-box volume (DOK 3).",
        ],
        footer="\u201cPencils up.  Exit ticket.  Savvas LQ Q5.  5 minutes.\u201d")

    add_phase_slide(prs, 7, total,
        "Exit Ticket  \u2014  Savvas Lesson Quiz Q5",
        "Exit Ticket", "DOK 2", 5, [
            "Savvas Lesson Quiz Q5:",
            "",
            "Find the zeros of  f(x) = \u2212x\u00b3 \u2212 2x\u00b2 + 7x \u2212 4.",
            "Then describe the behavior of the graph at each zero.",
            "",
            "Use multiplicity language:  \u201ccrosses\u201d / \u201ctouches\u201d / \u201cturns\u201d.",
            "",
            "Place your packet on the front desk as you leave.",
        ],
        footer="Packets to the front. Pack up. See you next class.",
        body_size=24)

    sw = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(sw, NAVY)
    add_text(sw, "Nice work today.",
             Inches(0.6), Inches(2.3), Inches(12), Inches(1.2),
             size=50, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(sw, [
        "Collect:  Exit Ticket  \u00b7  Do Now sheet  \u00b7  Practice page",
        "",
        "Day 4-5 preview:  real + complex zeros + storage-box modeling.",
    ],
        Inches(0.6), Inches(4.1), Inches(12), Inches(2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF),
        align=PP_ALIGN.CENTER)

    prs.save(path)


if __name__ == "__main__":
    build("Day_23_Slides.pptx")
    print("Built Day_23_Slides.pptx")
