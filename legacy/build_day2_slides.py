"""Generate Day_2_Slides.pptx \u2014 classroom projection deck for Day 2.

Design (revised): single-DOK3 spine.  7 phases.  The Reverse Engineering
slide is the driver and is intentionally information-dense so students can
read the page and work without waiting on the teacher.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x3C, 0x6C)
BLUE = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
DARK = RGBColor(0x11, 0x22, 0x33)
GRAY = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)


def add_background(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_badge(slide, text, left, top, *, color=NAVY, text_color=RGBColor(0xFF, 0xFF, 0xFF)):
    width = Inches(0.28 + 0.11 * len(text))
    height = Inches(0.35)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shp, width


def add_header(slide, phase_num, total, phase_title, framework, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_text(slide, phase_title, Inches(0.5), Inches(0.18),
             Inches(9.5), Inches(0.9),
             size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    meta = f"Phase {phase_num}/{total}  \u00b7  Algebra 2  \u00b7  Lesson 3-5  \u00b7  Day 2"
    add_text(slide, meta, Inches(0.5), Inches(0.82),
             Inches(9), Inches(0.35),
             size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    bx = Inches(10.3)
    by = Inches(0.4)
    if framework:
        shp, w = add_badge(slide, framework, bx, by, color=RGBColor(0x22, 0x33, 0x44))
        bx = bx + w + Inches(0.08)
    if dok and dok != "\u2014":
        add_badge(slide, dok, bx, by, color=BLUE)

    if minutes is not None:
        mt = f"{minutes} min"
        add_text(slide, mt, Inches(10.3), Inches(0.85),
                 Inches(2.5), Inches(0.35),
                 size=14, bold=True,
                 color=RGBColor(0xCC, 0xDD, 0xEE),
                 align=PP_ALIGN.LEFT)


def add_phase_slide(prs, phase_num, total, phase_title, framework, dok,
                    minutes, body_lines, footer=None, body_size=24):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, phase_num, total, phase_title, framework, dok, minutes)
    add_text(slide, body_lines,
             Inches(0.6), Inches(1.6),
             Inches(12.1), Inches(5.3),
             size=body_size, color=DARK)
    if footer:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     0, Inches(6.85), SLIDE_W, Inches(0.65))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        add_text(slide, footer, Inches(0.5), Inches(6.92),
                 Inches(12.3), Inches(0.5),
                 size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    return slide


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---------- Slide 1: Title ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "ALGEBRA 2  \u00b7  LESSON 3-5",
             Inches(0.6), Inches(1.3), Inches(12), Inches(0.7),
             size=24, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, "Day 2  \u2014  Graphing from Factored Form",
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.4),
             size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, [
        "Essential Question:",
        "How do the zeros \u2014 and a single point on the graph \u2014",
        "determine the polynomial\u2019s equation?",
        "",
        "Materials: laptop (Desmos), blank paper, pencil, Blooket code.",
    ],
        Inches(0.6), Inches(3.6), Inches(12), Inches(3.2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_text(s, "57-minute lesson  \u00b7  one DOK-3 driver",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             size=16, color=RGBColor(0x88, 0xAA, 0xCC))

    total = 7

    # ---------- Slide 2: Do Now A ----------
    add_phase_slide(prs, 1, total,
        "Do Now  \u2014  Sign Flip Prediction",
        "Do Now A", "DOK 2", 5, [
            "On the Do Now sheet in front of you:",
            "",
            "Compare   f(x) = x(x \u2212 4)(x + 3)",
            "  vs.      g(x) = \u2212x(x \u2212 4)(x + 3)",
            "",
            "1.  What STAYS THE SAME?  What CHANGES?  Be specific.",
            "",
            "2.  Finish:  \u201cThe zeros stay the same because ___,",
            "                  and the graph changes because ___.\u201d",
            "",
            "Silent  \u00b7  Pencil only  \u00b7  No Desmos  \u00b7  No packet yet",
        ],
        footer="When everyone\u2019s back:  Blooket code on the screen.")

    # ---------- Slide 3: Blooket Login ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, 2, total, "Blooket  \u2014  Log In", "Do Now B", "\u2014", 2)
    add_text(s, "Blooket code:",
             Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
             size=28, color=GRAY)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.5), Inches(2.6), Inches(10.3), Inches(3))
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = NAVY
    box.line.width = Pt(3)
    add_text(s, "( write the code here or paste a screenshot )",
             Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.5),
             size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "Chromebooks open  \u00b7  Type the code  \u00b7  Enter a nickname",
             Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(6.85), SLIDE_W, Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(s, "2 minutes to log in.  Game waits for no one.",
             Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.5),
             size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # ---------- Slide 4: Blooket Rule Recall ----------
    add_phase_slide(prs, 3, total,
        "Blooket  \u2014  Rule Recall",
        "Do Now C", "DOK 1", 7, [
            "Pure rule recall today.  The rules you\u2019ll need in 20 minutes:",
            "",
            "  \u2022  Zero Product Property  \u2014  zero of f \u21d4 factor = 0",
            "  \u2022  Sign direction  \u2014  zero at x = \u22122 gives (x + 2)",
            "  \u2022  x\u00b2 + 9 (or any x\u00b2 + positive) has  NO real zeros",
            "  \u2022  Highest power \u2192 end behavior",
            "  \u2022  Zeros of \u2212f(x) are the SAME as zeros of f(x)",
            "",
            "Play fast.  Teacher is watching the dashboard.",
        ],
        footer="Close game at 0.  \u201cPull out your Do Now.  Packet too.\u201d")

    # ---------- Slide 5: Launch \u2014 Sign Flip Synthesis ----------
    add_phase_slide(prs, 4, total,
        "Launch  \u2014  Sign Flip Synthesis",
        "Launch", "DOK 2", 8, [
            "Out:  Do Now sheet  +  Student packet.",
            "",
            "1.  Turn and Talk with your partner:",
            "      What did you predict would STAY THE SAME?",
            "      What did you predict would CHANGE?",
            "",
            "2.  Class share  +  Desmos check.",
            "",
            "3.  Sentence frame (say it aloud, write it later in Step 5):",
            "      \u201cThe zeros stay the same because ___.",
            "       The graph changes because ___.\u201d",
        ],
        footer="\u201cPens down.  Now: BACKWARD.  Zeros + a point \u2192 equation.\u201d")

    # ---------- Slide 6: Reverse Engineering (THE DOK-3 DRIVER) ----------
    add_phase_slide(prs, 5, total,
        "Explore  \u2014  Build the Equation (Reverse)",
        "Explore", "DOK 3", 25, [
            "Write the equation of a polynomial with",
            "    zeros at  x = \u22122,  x = 1,  x = 4",
            "    passing through the point  (0, \u22128).",
            "",
            "EVERYTHING YOU NEED IS ON YOUR PACKET PAGE:",
            "  Rule 1:  zero at c  \u2192  factor (x \u2212 c)",
            "  Rule 2:  f(x) = a \u00b7 (factor)(factor)(factor) \u2014 use the point to find a.",
            "",
            "  Step 1  Zeros \u2192 factors.",
            "  Step 2  Plug in (0, \u22128) to solve for a.",
            "  Step 3  Final equation.    Step 4  Verify in Desmos.",
            "  Step 5  Justify to partner (use the Do Now sentence frame).",
        ],
        footer="Teacher circulates.  Prompts only \u2014 answer questions with questions.",
        body_size=22)

    # ---------- Slide 7: Share / Summary ----------
    add_phase_slide(prs, 6, total,
        "Share / Summary",
        "Share/Summary", "DOK 2", 5, [
            "Essential Question callback:",
            "How do the zeros \u2014 and a single point on the graph \u2014",
            "determine the polynomial\u2019s equation?",
            "",
            "Self-rating on your packet:   \u2713  /  partly  /  not yet",
            "  1. Zeros \u2192 factors",
            "  2. Point \u2192 leading number a",
            "  3. Predicting what changes under \u2212f(x)",
            "",
            "Preview Day 3:  what happens when a factor appears twice?",
        ],
        footer="\u201cPencils up.  Summary exit ticket.  5 minutes.\u201d")

    # ---------- Slide 8: Exit Ticket \u2014 SUMMARY ----------
    add_phase_slide(prs, 7, total,
        "Exit Ticket  \u2014  Summary",
        "Exit Ticket", "DOK 1\u20132", 5, [
            "1.  When f(x) becomes \u2212f(x), the zeros ______",
            "     and the graph ______ .",
            "",
            "2.  To find the leading number  a , I plug in a given ______",
            "     and solve for  a .",
            "",
            "3.  One sentence \u2014 the most important thing I learned today.",
            "",
            "Place your packet on the front desk as you leave.",
        ],
        footer="Packets to the front. Pack up. See you next class.",
        body_size=26)

    # ---------- Slide 9: Wrap-up ----------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "Nice work today.",
             Inches(0.6), Inches(2.3), Inches(12), Inches(1.2),
             size=52, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(s, [
        "Collect:  Exit Ticket  \u00b7  Do Now sheet  \u00b7  Reverse Engineering page",
        "",
        "Day 3 preview:  what happens when a factor repeats?",
    ],
        Inches(0.6), Inches(4.1), Inches(12), Inches(2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF),
        align=PP_ALIGN.CENTER)

    prs.save(path)


if __name__ == "__main__":
    build("Day_2_Slides.pptx")
    print("Built Day_2_Slides.pptx")
