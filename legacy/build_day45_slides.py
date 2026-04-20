"""Day_45_Slides.pptx \u2014 Combined Day 4-5 projection deck (Savvas-only).

Eight phase slides + title + wrap.  DOK-3 capstone is Savvas Practice #27
(storage box volume), the Savvas-declared DOK-3 anchor for this lesson.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x3C, 0x6C)
BLUE = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
DARK = RGBColor(0x11, 0x22, 0x33)
GRAY = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)


def add_background(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_badge(slide, text, left, top, *, color=NAVY,
              text_color=RGBColor(0xFF, 0xFF, 0xFF)):
    width = Inches(0.28 + 0.11 * len(text))
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, Inches(0.35))
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shp, width


def add_header(slide, phase_num, total, phase_title, framework, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_text(slide, phase_title, Inches(0.5), Inches(0.18),
             Inches(9.5), Inches(0.9),
             size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    meta = f"Phase {phase_num}/{total}  \u00b7  Algebra 2  \u00b7  Lesson 3-5  \u00b7  Day 4-5"
    add_text(slide, meta, Inches(0.5), Inches(0.82),
             Inches(9), Inches(0.35),
             size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    bx = Inches(10.3)
    by = Inches(0.4)
    if framework:
        shp, w = add_badge(slide, framework, bx, by,
                           color=RGBColor(0x22, 0x33, 0x44))
        bx = bx + w + Inches(0.08)
    if dok and dok != "\u2014":
        add_badge(slide, dok, bx, by, color=BLUE)

    if minutes is not None:
        add_text(slide, f"{minutes} min", Inches(10.3), Inches(0.85),
                 Inches(2.5), Inches(0.35),
                 size=14, bold=True,
                 color=RGBColor(0xCC, 0xDD, 0xEE))


def add_phase_slide(prs, phase_num, total, phase_title, framework, dok,
                    minutes, body_lines, footer=None, body_size=22):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, phase_num, total, phase_title, framework, dok, minutes)
    add_text(slide, body_lines,
             Inches(0.6), Inches(1.6),
             Inches(12.1), Inches(5.3),
             size=body_size, color=DARK)
    if footer:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     0, Inches(6.85), SLIDE_W, Inches(0.65))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        add_text(slide, footer, Inches(0.5), Inches(6.92),
                 Inches(12.3), Inches(0.5),
                 size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    return slide


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "ALGEBRA 2  \u00b7  LESSON 3-5",
             Inches(0.6), Inches(1.3), Inches(12), Inches(0.7),
             size=22, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, "Combined Day 4-5  \u2014  Real + Complex Zeros + Modeling",
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.4),
             size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, [
        "Essential Question:",
        "How do synthetic division and the quadratic formula let you find",
        "EVERY zero of a polynomial \u2014 and what do those zeros mean when",
        "the polynomial models a real-world quantity like volume?",
        "",
        "DOK-3 capstone: Savvas Practice #27 (storage box volume).",
    ],
        Inches(0.6), Inches(3.6), Inches(12), Inches(3.2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_text(s, "55-min lesson  \u00b7  Savvas-declared DOK 3",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             size=16, color=RGBColor(0x88, 0xAA, 0xCC))

    total = 8

    add_phase_slide(prs, 1, total,
        "Do Now  \u2014  Zeros \u21d4 Factors",
        "Do Now A", "DOK 1", 5, [
            "SAVVAS PRACTICE #28  \u2014  Vocabulary",
            "",
            "Complete each statement so it means the same as",
            "\u201c4 is a zero of the function\u201d:",
            "  (1)  The graph crosses the _______ at 4.",
            "  (2)  _______ is a factor of the polynomial.",
            "",
            "Bridge:  If \u22122 is a zero, one factor is _______.",
            "         A degree-3 polynomial with only 1 real zero",
            "         must have _______ complex zeros.",
            "",
            "Silent  \u00b7  Pencil only  \u00b7  No Desmos",
        ],
        footer="When everyone\u2019s in:  Blooket code on the screen.")

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s2)
    add_header(s2, 2, total, "Blooket  \u2014  Log In", "Do Now B", "\u2014", 2)
    add_text(s2, "Blooket code:",
             Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
             size=28, color=GRAY)
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.5), Inches(2.6), Inches(10.3), Inches(3))
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = NAVY
    box.line.width = Pt(3)
    add_text(s2, "( write the code here )",
             Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.5),
             size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s2, "Chromebooks open  \u00b7  Type code  \u00b7  Nickname",
             Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.85),
                              SLIDE_W, Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(s2, "2 minutes to log in. Game waits for no one.",
             Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.5),
             size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    add_phase_slide(prs, 3, total,
        "Blooket  \u2014  Rule Recall",
        "Do Now C", "DOK 1", 7, [
            "Rules today\u2019s work depends on:",
            "",
            "  \u2022  Degree n  \u21d2  n total zeros (including complex)",
            "  \u2022  Complex zeros come in CONJUGATE PAIRS  (a + bi with a \u2212 bi)",
            "  \u2022  Factor theorem:  x = c  \u21d4  (x \u2212 c)",
            "  \u2022  Synthetic division procedure",
            "  \u2022  Quadratic formula on the reduced quadratic",
            "",
            "Play fast.  Teacher is watching the dashboard.",
        ],
        footer="Close game at 0.  \u201cPull out your Do Now. Packet too.\u201d")

    add_phase_slide(prs, 4, total,
        "Launch  \u2014  Savvas Example 3",
        "Launch", "DOK 2", 10, [
            "Graph in Desmos:  f(x) = x\u00b3 \u2212 2x\u00b2 \u2212 2x \u2212 3",
            "",
            "  Step 1.  How many times does it cross?  One real zero:  x = 3.",
            "  Step 2.  Degree 3 \u2192 3 total zeros.  2 unaccounted for.",
            "  Step 3.  Synthetic division by (x \u2212 3)  \u2192  Q(x) = x\u00b2 + x + 1",
            "  Step 4.  Quadratic formula  \u2192  x = (\u22121 \u00b1 i\u221a3) / 2",
            "",
            "THE RULE:  non-real zeros come in CONJUGATE PAIRS.",
            "THE SHORTCUT:  degree n  \u21d2  n total zeros.",
        ],
        footer="\u201cPackets to Practice.  Try It 3a and 3b.  10 minutes.\u201d")

    add_phase_slide(prs, 5, total,
        "Practice  \u2014  Savvas Try It 3a + 3b",
        "Practice", "DOK 2", 10, [
            "For each, find ALL real and complex zeros:",
            "",
            "  A.  f(x) = 2x\u00b3 \u2212 8x\u00b2 + 9x           [Try It 3a]",
            "  B.  f(x) = x\u2074 \u2212 3x\u00b2 \u2212 4               [Try It 3b]",
            "",
            "Show your synthetic-division work.",
            "Write complex zeros using i.",
            "Verify:  total zero count = degree.",
        ],
        footer="\u201cPacket flip.  Storage box next.  Read the rules on the page.\u201d")

    add_phase_slide(prs, 6, total,
        "Storage Box  \u2014  the DOK-3 driver",
        "Explore", "DOK 3", 15, [
            "SAVVAS PRACTICE #27  \u2014  Model With Mathematics",
            "",
            "Height is LESS than both length and width.",
            "Volume:  f(x) = x\u00b3 + 2x\u00b2 \u2212 3x  where x = width (ft).",
            "",
            "  (a)  Factor f(x).",
            "  (b)  Find the zeros.",
            "  (c)  Match factors to dimensions (use \u201cheight is smaller\u201d).",
            "  (d)  Find dimensions when V = 10 ft\u00b3.",
            "",
            "All rules are on your packet page.  Teacher circulates with prompts only.",
        ],
        footer="Stuck after 10 min?  Hint card: \u201cFactor first.  Which factor is "
               "smaller at x = 2?\u201d",
        body_size=20)

    add_phase_slide(prs, 7, total,
        "Share / Summary",
        "Share/Summary", "DOK 2", 3, [
            "Essential Question:",
            "How do synthetic division and the quadratic formula",
            "let us find EVERY zero?",
            "",
            "Self-rating on your packet:   \u2713  /  partly  /  not yet",
            "",
            "Preview Day 6:  polynomial inequalities \u2014 where is f(x) > 0?",
        ],
        footer="\u201cPencils up. Exit ticket. Savvas #17. 5 minutes.\u201d")

    add_phase_slide(prs, 8, total,
        "Exit Ticket  \u2014  Savvas Practice #17",
        "Exit Ticket", "DOK 2", 5, [
            "Savvas Practice #17:",
            "",
            "Find all real and complex zeros of the polynomial function",
            "shown in the graph.",
            "",
            "Then verify: total zero count = degree.",
            "",
            "Place your packet on the front desk as you leave.",
        ],
        footer="Packets to the front. Pack up. See you next class.",
        body_size=24)

    sw = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(sw, NAVY)
    add_text(sw, "Nice work today.",
             Inches(0.6), Inches(2.3), Inches(12), Inches(1.2),
             size=50, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(sw, [
        "Collect:  Exit Ticket  \u00b7  Do Now  \u00b7  Launch + Practice + Storage Box pages",
        "",
        "Day 6 preview:  polynomial inequalities.",
    ],
        Inches(0.6), Inches(4.1), Inches(12), Inches(2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF),
        align=PP_ALIGN.CENTER)

    prs.save(path)


if __name__ == "__main__":
    build("Day_45_Slides.pptx")
    print("Built Day_45_Slides.pptx")
