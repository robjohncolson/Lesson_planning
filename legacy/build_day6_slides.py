"""Day_6_Slides.pptx — Combined Day 6 projection deck (Savvas-only).

Ten slides: title + one per phase + wrap.
Every work item traces to a Savvas bank ID. Sign-chart method walk-through
on the Launch slide so students scanning the screen see the same steps they
are working on the packet.
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
TEAL   = RGBColor(0x00, 0x97, 0xA7)


def add_background(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.1)
    tf.margin_right = Inches(0.1)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name  = font
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
    return tb


def add_badge(slide, text, left, top, *, color=NAVY,
              text_color=RGBColor(0xFF, 0xFF, 0xFF)):
    width = Inches(0.28 + 0.11 * len(text))
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, Inches(0.35))
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left  = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name  = "Calibri"
    run.font.size  = Pt(11)
    run.font.bold  = True
    run.font.color.rgb = text_color
    return shp, width


def add_header(slide, phase_num, total, phase_title, framework, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_text(slide, phase_title, Inches(0.5), Inches(0.18),
             Inches(9.5), Inches(0.9),
             size=34, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    meta = f"Phase {phase_num}/{total}  \u00b7  Algebra 2  \u00b7  Lesson 3-5  \u00b7  Day 6"
    add_text(slide, meta, Inches(0.5), Inches(0.82),
             Inches(9), Inches(0.35),
             size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    bx = Inches(10.3)
    by = Inches(0.4)
    if framework:
        shp, w = add_badge(slide, framework, bx, by,
                           color=RGBColor(0x22, 0x33, 0x44))
        bx = bx + w + Inches(0.08)
    if dok and dok != "\u2014":
        add_badge(slide, dok, bx, by, color=BLUE)

    if minutes is not None:
        add_text(slide, f"{minutes} min", Inches(10.3), Inches(0.85),
                 Inches(2.5), Inches(0.35),
                 size=14, bold=True,
                 color=RGBColor(0xCC, 0xDD, 0xEE))


def add_phase_slide(prs, phase_num, total, phase_title, framework, dok,
                    minutes, body_lines, footer=None, body_size=22):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, phase_num, total, phase_title, framework, dok, minutes)
    add_text(slide, body_lines,
             Inches(0.6), Inches(1.6),
             Inches(12.1), Inches(5.3),
             size=body_size, color=DARK)
    if footer:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     0, Inches(6.85), SLIDE_W, Inches(0.65))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        add_text(slide, footer, Inches(0.5), Inches(6.92),
                 Inches(12.3), Inches(0.5),
                 size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    return slide


def build(path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Title ────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "ALGEBRA 2  \u00b7  LESSON 3-5",
             Inches(0.6), Inches(1.3), Inches(12), Inches(0.7),
             size=22, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, "Combined Day 6  \u2014  Polynomial Inequalities",
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.4),
             size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, [
        "Essential Question:",
        "How do the zeros of a polynomial divide the number line,",
        "and how does that structure let you solve any polynomial inequality?",
        "",
        "Every work item today comes from Savvas.",
    ],
        Inches(0.6), Inches(3.6), Inches(12), Inches(3.2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_text(s, "55-min lesson  \u00b7  Savvas-only (no fabricated items)",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             size=16, color=RGBColor(0x88, 0xAA, 0xCC))

    total = 9

    # ── Slide 2: Do Now A ─────────────────────────────────────────────
    add_phase_slide(prs, 1, total,
        "Do Now  \u2014  Solve the Equation (Bridge)",
        "Do Now A", "DOK 1", 5, [
            "On the Do Now sheet  \u2014  SAVVAS PRACTICE #21:",
            "",
            "\u221a5x\u2074 + 4x\u00b2 \u2212 12x = \u22126x\u2074 + 3x\u00b3",
            "",
            "Find ALL solutions \u2014 real AND complex.  Show your work.",
            "",
            "Silent  \u00b7  Pencil only  \u00b7  No Desmos",
            "",
            "Hold on to your zeros.  We\u2019re going to put them on a number line.",
        ],
        footer="When everyone\u2019s in:  Blooket code on the next slide.")

    # ── Slide 3: Blooket Login ────────────────────────────────────────
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s2)
    add_header(s2, 2, total, "Blooket  \u2014  Log In", "Do Now B", "\u2014", 2)
    add_text(s2, "Blooket code:",
             Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
             size=28, color=GRAY)
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.5), Inches(2.6), Inches(10.3), Inches(3))
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = NAVY
    box.line.width = Pt(3)
    add_text(s2, "( write the code here )",
             Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.5),
             size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s2, "Chromebooks open  \u00b7  Type the code  \u00b7  Enter nickname",
             Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.85),
                              SLIDE_W, Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(s2, "2 minutes to log in. Game waits for no one.",
             Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.5),
             size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # ── Slide 4: Blooket Rule Recall ──────────────────────────────────
    add_phase_slide(prs, 3, total,
        "Blooket  \u2014  Rule Recall",
        "Do Now C", "DOK 1", 6, [
            "Pure rule recall.  The rules you\u2019ll need in 30 minutes:",
            "",
            "  \u2022  Zero \u21d4 factor = 0  (Zero Product Property)",
            "  \u2022  Sign of f can only change at a zero of ODD multiplicity",
            "  \u2022  Even-multiplicity zero = same sign on both sides",
            "  \u2022  Continuous polynomial \u2192 sign is constant between zeros",
            "  \u2022  GCF first; use Quadratic Formula if the factor doesn\u2019t split",
            "  \u2022  Strict inequality (<, >) \u2192 zeros NOT included in solution",
            "",
            "Play fast.  Teacher is watching the dashboard.",
        ],
        footer="Close game at 0.  Blooket Wrap: candy + lowest rule aloud.")

    # ── Slide 5: Launch — Sign-Chart Method ───────────────────────────
    add_phase_slide(prs, 4, total,
        "Launch  \u2014  Sign-Chart Method",
        "Launch", "DOK 2", 12, [
            "Savvas Example 6 / Try It 6a:  Solve  2x\u00b3 + 12x\u00b2 + 12x < 0",
            "",
            "  Step 1  Move everything to one side.  Factor out the GCF.",
            "          2x(x\u00b2 + 6x + 6) < 0",
            "",
            "  Step 2  Find zeros.  (Quadratic Formula on x\u00b2 + 6x + 6.)",
            "          x = 0,   x = \u22123 \u2212 \u221a3,   x = \u22123 + \u221a3",
            "",
            "  Step 3  Plot zeros on number line.  Label the intervals.",
            "",
            "  Step 4  Test one point per interval. Record sign (+) or (\u2212).",
            "",
            "  Step 5  Shade intervals where sign = (\u2212).  Read the solution.",
        ],
        footer="\u201cPackets to Practice \u2014 20 minutes. A required; pick two of B/C/D.\u201d",
        body_size=20)

    # ── Slide 6: Sign-Chart Rules (keep visible during Practice) ──────
    add_phase_slide(prs, 4, total,
        "Sign-Chart Rules  \u2014  Reference",
        "Launch", "DOK 2", None, [
            "RULE 1   Move to one side \u2192 f(x) < 0  or  f(x) > 0.",
            "RULE 2   Factor completely.  Find all real zeros.",
            "RULE 3   Zeros divide the number line into intervals.",
            "RULE 4   Test ONE point per interval.  Sign is constant on the interval.",
            "RULE 5   Sign changes only at zeros of ODD multiplicity.",
            "          Even-multiplicity zeros: sign does NOT change.",
            "RULE 6   Shade where the sign matches the inequality.",
            "          Strict (<, >) \u2192 zeros excluded.   Non-strict (\u2264, \u2265) \u2192 included.",
            "",
            "Sentence frame:  \u201cThe zeros are ___, so the sign changes at ___,",
            "                   and the solution is ___.\u201d",
        ],
        footer="Keep this on screen during Practice.",
        body_size=21)

    # ── Slide 7: Practice ─────────────────────────────────────────────
    add_phase_slide(prs, 5, total,
        "Practice  \u2014  Savvas Try It 6b + #22/#23/#24",
        "Practice", "DOK 2", 20, [
            "Sign chart for each.  Required: A.   Pick two of B, C, D.",
            "",
            "  A.  (x\u00b2 \u2212 1)(x\u00b2 \u2212 x \u2212 6) > 0   [Try It 6b]",
            "  B.  x\u00b3 \u2212 9x > 0               [Practice #22]",
            "  C.  0 > 4x\u00b3 + 8x\u00b2 \u2212 x \u2212 2      [Practice #23]",
            "  D.  64x\u00b2 > \u22124x\u00b3 \u2212 x \u2212 16       [Practice #24]",
            "",
            "For each item:",
            "  1. Factor  \u2192  2. Find zeros  \u2192  3. Plot on number line",
            "  4. Test points  \u2192  5. Shade  \u2192  6. Write solution",
        ],
        footer="Teacher circulates.  Prompts only \u2014 answer questions with questions.")

    # ── Slide 8: Share / Summary ───────────────────────────────────────
    add_phase_slide(prs, 6, total,
        "Share / Summary",
        "Share/Summary", "DOK 2", 4, [
            "Essential Question:",
            "How do the zeros of a polynomial divide the number line?",
            "How does that let you solve any polynomial inequality?",
            "",
            "Self-rating on your packet:   \u2713  /  partly  /  not yet",
            "  1.  I can factor and identify all real zeros.",
            "  2.  I can build a sign chart with correct intervals.",
            "  3.  I can read the solution set from the sign chart.",
            "",
            "One-sentence synthesis:  \u201cThe key idea today was ___.\u201d",
            "",
            "Preview Day 7:  synthesis + Savvas Performance Task #30 (Venetta, DOK 3).",
        ],
        footer="\u201cPencils up.  Exit ticket.  LQ Q3.  5 minutes.\u201d")

    # ── Slide 9: Exit Ticket ───────────────────────────────────────────
    add_phase_slide(prs, 7, total,
        "Exit Ticket  \u2014  Savvas Lesson Quiz Q3",
        "Exit Ticket", "DOK 2", 5, [
            "Savvas Lesson Quiz Q3:",
            "",
            "What values of x are solutions of the inequality  x\u00b3 \u2212 4x > 0?",
            "",
            "Show your sign chart.  Write your solution in interval notation.",
            "",
            "Desmos: YES, to verify zeros only.  Sign chart = your work.",
            "",
            "Place your packet on the front desk as you leave.",
        ],
        footer="Packets to the front.  Pack up.  See you next class.",
        body_size=24)

    # ── Slide 10: Wrap / Class End ────────────────────────────────────
    sw = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(sw, NAVY)
    add_text(sw, "Nice work today.",
             Inches(0.6), Inches(2.3), Inches(12), Inches(1.2),
             size=50, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(sw, [
        "Collect:  Exit Ticket  \u00b7  Do Now sheet  \u00b7  Practice page",
        "",
        "Day 7 preview:  Synthesis + Savvas Performance Task #30 (DOK 3).",
    ],
        Inches(0.6), Inches(4.1), Inches(12), Inches(2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF),
        align=PP_ALIGN.CENTER)

    prs.save(path)


if __name__ == "__main__":
    build("Day_6_Slides.pptx")
    print("Built Day_6_Slides.pptx")
