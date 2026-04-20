"""Generate Day_3_Slides.pptx - classroom projection deck for Day 3 (Multiplicity).

One slide per phase with [Framework] + [DOK] badges. Built for 55-min block.
Payoff day: we finally NAME the phenomenon students saw at x = 4 yesterday.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x3C, 0x6C)
BLUE = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
DARK = RGBColor(0x11, 0x22, 0x33)
GRAY = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)


def add_background(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_badge(slide, text, left, top, *, color=NAVY, text_color=RGBColor(0xFF, 0xFF, 0xFF)):
    width = Inches(0.28 + 0.11 * len(text))
    height = Inches(0.35)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shp, width


def add_header(slide, phase_num, total, phase_title, framework, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_text(slide, phase_title, Inches(0.5), Inches(0.18),
             Inches(9.5), Inches(0.9),
             size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    meta = f"Phase {phase_num}/{total}  \u00b7  Algebra 2  \u00b7  Lesson 3-5  \u00b7  Day 3"
    add_text(slide, meta, Inches(0.5), Inches(0.82),
             Inches(9), Inches(0.35),
             size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    bx = Inches(10.3)
    by = Inches(0.4)
    if framework:
        shp, w = add_badge(slide, framework, bx, by, color=RGBColor(0x22, 0x33, 0x44))
        bx = bx + w + Inches(0.08)
    if dok and dok != "\u2014":
        shp, w = add_badge(slide, dok, bx, by, color=BLUE)

    if minutes is not None:
        mt = f"{minutes} min"
        add_text(slide, mt, Inches(10.3), Inches(0.85),
                 Inches(2.5), Inches(0.35),
                 size=14, bold=True,
                 color=RGBColor(0xCC, 0xDD, 0xEE),
                 align=PP_ALIGN.LEFT)


def add_phase_slide(prs, phase_num, total, phase_title, framework, dok,
                    minutes, body_lines, footer=None, body_size=24):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, phase_num, total, phase_title, framework, dok, minutes)
    add_text(slide, body_lines,
             Inches(0.6), Inches(1.6),
             Inches(12.1), Inches(5.3),
             size=body_size, color=DARK)
    if footer:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     0, Inches(6.85), SLIDE_W, Inches(0.65))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        add_text(slide, footer, Inches(0.5), Inches(6.92),
                 Inches(12.3), Inches(0.5),
                 size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    return slide


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "ALGEBRA 2  \u00b7  LESSON 3-5",
             Inches(0.6), Inches(1.3), Inches(12), Inches(0.7),
             size=24, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, "Day 3  \u2014  The Magic of Multiplicity",
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.4),
             size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, [
        "Essential Question:",
        "What does a repeated factor do to the graph of a polynomial,",
        "and how can you predict the behavior at every zero just by looking",
        "at the factored form?",
        "",
        "Day 2: zeros + a point \u2192 equation (you worked BACKWARD).",
        "Today: when a factor appears more than once, the graph does",
        "something different at that zero. We\u2019ll name it.",
    ],
        Inches(0.6), Inches(3.6), Inches(12), Inches(3.4),
        size=20, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_text(s, "55-min lesson  \u00b7  one DOK-3 driver",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             size=16, color=RGBColor(0x88, 0xAA, 0xCC))

    total = 7

    # 1. Do Now
    add_phase_slide(prs, 1, total,
        "Do Now  \u2014  Touch or Cross?",
        "Do Now A", "DOK 2", 5, [
            "On the Do Now sheet in front of you \u2014",
            "for EACH zero, predict: CROSS or TOUCH?",
            "",
            "1.  f(x) = (x + 2)(x \u2212 3)",
            "2.  h(x) = (x \u2212 1)\u00b2 (x + 5)",
            "",
            "3.  Pattern guess:  EVEN factor count \u2192 ?    ODD factor count \u2192 ?",
            "",
            "Silent  \u00b7  Pencil only  \u00b7  No Desmos  \u00b7  No packet",
        ],
        footer="When everyone\u2019s in:  Blooket code coming up.")

    # 2. Blooket login
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, 2, total, "Blooket  \u2014  Log In", "Do Now B", "\u2014", 2)
    add_text(s, "Blooket code:",
             Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
             size=28, color=GRAY)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.5), Inches(2.6), Inches(10.3), Inches(3))
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = NAVY
    box.line.width = Pt(3)
    add_text(s, "( Day 3 Multiplicity Blooket code here )",
             Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.5),
             size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "Chromebooks open \u00b7 Type the code \u00b7 Enter a nickname",
             Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(6.85), SLIDE_W, Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(s, "2 minutes to log in. Game waits for no one.",
             Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.5),
             size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # 3. Blooket game
    add_phase_slide(prs, 3, total,
        "Blooket  \u2014  Rule Recall",
        "Do Now C", "DOK 1", 7, [
            "Pure rule recall today. The rules you\u2019ll need in 30 minutes:",
            "",
            "  \u2022  Multiplicity = exponent on the factor",
            "  \u2022  EVEN multiplicity \u2192 graph TOUCHES and turns",
            "  \u2022  ODD multiplicity \u2192 graph CROSSES",
            "  \u2022  Day 2 recall: zero \u21d4 factor = 0  /  zero at c gives (x \u2212 c)",
            "",
            "Play fast.  Teacher is watching the dashboard.",
        ],
        footer="Close game at 0.  \u201cPull out your Do Now. Packet too.\u201d")

    # 4. Launch \u2014 Name the Pattern via (x-1)\u207f
    add_phase_slide(prs, 4, total,
        "Launch  \u2014  Name the Pattern",
        "Launch", "DOK 2", 12, [
            "1.  Pull out your Do Now.  Cold-call the pattern in #3.",
            "",
            "2.  In Desmos \u2014 type each, watch what happens at x = 1:",
            "      (a)  y = (x \u2212 1)\u00b9        (b)  y = (x \u2212 1)\u00b2",
            "      (c)  y = (x \u2212 1)\u00b3        (d)  y = (x \u2212 1)\u2074",
            "",
            "3.  THE WORD:  the number of times a factor appears in the factored",
            "    form is its  MULTIPLICITY.   (x \u2212 1)\u00b2  has multiplicity 2 at x = 1.",
            "",
            "4.  Fill the T-chart on your packet:  ODD  vs.  EVEN.",
        ],
        footer="\u201cPackets to Practice. Three quick items.\u201d",
        body_size=22)

    # 5. Practice \u2014 Try It A-C
    add_phase_slide(prs, 5, total,
        "Practice  \u2014  Try It A\u2013C",
        "Practice", "DOK 2", 10, [
            "For each polynomial:  list zeros, multiplicity, and behavior.",
            "",
            "  A.  f(x) = (x + 3)(x \u2212 2)\u00b2",
            "  B.  g(x) = x\u00b2 (x \u2212 5)\u00b3",
            "  C.  h(x) = (x + 1)\u2074 (x \u2212 4)",
            "",
            "Frame:  \u201cThe multiplicity is ____, so the graph ____ at x = ____.\u201d",
            "",
            "Fast finishers \u2014 challenge on the board.",
        ],
        footer="\u201cPacket flip. Tonya\u2019s graph next. Read the rules card on the page.\u201d")

    # 6. Tonya \u2014 the DOK-3 driver
    add_phase_slide(prs, 6, total,
        "Explore  \u2014  Tonya\u2019s Graph (Error Analysis)",
        "Explore", "DOK 3", 15, [
            "Tonya sketches   f(x) = (x + 1)\u00b2 (x \u2212 3).",
            "She draws the graph CROSSING at x = \u22121 AND CROSSING at x = 3.",
            "",
            "EVERYTHING YOU NEED IS ON YOUR PAGE  \u2014  three rules + 5 steps.",
            "",
            "  Step 1  What did Tonya get RIGHT?",
            "  Step 2  What did she get WRONG?  Use \u201cmultiplicity\u201d.",
            "  Step 3  Defend the correct behavior at x = \u22121 (sentence frame).",
            "  Step 4  Verify in Desmos.    Step 5  Explain to your partner.",
        ],
        footer="Teacher circulates.  Prompts only \u2014 \u201cWhat\u2019s the exponent on (x + 1)?\u201d",
        body_size=22)

    # 7. Summary Exit Ticket
    add_phase_slide(prs, 7, total,
        "Exit Ticket  \u2014  Summary",
        "Exit Ticket", "DOK 1\u20132", 4, [
            "1.  MULTIPLICITY = ______ on the factor.",
            "",
            "2.  EVEN \u2192 ______ .    ODD \u2192 ______ .",
            "",
            "3.  In  p(x) = (x + 2)\u00b3 (x \u2212 4):",
            "      x = \u22122 \u2192 mult ____, behavior ______________",
            "      x =  4 \u2192 mult ____, behavior ______________",
            "",
            "4.  One sentence: most important thing I learned today.",
        ],
        footer="Packets to the front. Pack up. See you next class.",
        body_size=24)

    # 8. Wrap-up
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "Nice work.",
             Inches(0.6), Inches(2.3), Inches(12), Inches(1.2),
             size=52, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(s, [
        "Collect:  Exit Ticket  \u00b7  Do Now sheet  \u00b7  Tonya error analysis",
        "",
        "Day 4 preview:  real vs. complex zeros \u2014 what if we can\u2019t factor over the reals?",
    ],
        Inches(0.6), Inches(4.1), Inches(12), Inches(2),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF),
        align=PP_ALIGN.CENTER)

    prs.save(path)


if __name__ == "__main__":
    build("Day_3_Slides.pptx")
    print("Built Day_3_Slides.pptx")
