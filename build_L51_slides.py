"""L51_P1_Slides.pptx — projection deck for Lesson 5-1 (single-period).

Single deck: nth Roots, Rational Exponents + DOK-3 Cylinder Capstone (Topic 5 Assessment Item 1A Prep).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, period_label, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, f"Algebra 2  ·  Lesson 5-1  ·  {period_label}",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
         size=28, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Klimsara-adapted · Student-centered · No Blooket",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


RULES_TEXT = (
    "x^(1/n) = ⁿ√x.  "
    "x^(m/n) = (ⁿ√x)^m.  "
    "Rationalize ⁿ√x denom: × by ⁿ√(x^(n-1)).  "
    "Even n: negative radicand undefined in ℝ."
)


# ── P1 deck (single period) ───────────────────────────────────────────────────

def build_p1(path="L51_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(
        prs,
        "Lesson 5-1 · Quick — nth Roots, Rational Exponents + ⭐ Cylinder DOK-3 (Item 1A Prep)",
        "nth Roots + Rational Exponents + DOK-3 Cylinder Capstone",
    )

    # ── Do Now ────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Do Now — Explore & Reason",
           framework_tag="exploration", dok=3, minutes=5)
    card(s, "5-1-savvas-model-discuss-lesson-5-1-launch · Explore & Reason: y=x² graph",
         ["Look at the graph of y = x² on the board.",
          "",
          "A) What input gives output 9? Output 25? Output 2?",
          "B) Is there an input that gives output −1? Explain.",
          "C) Connect: what does the graph tell us about ROOTS?"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    # ── Launch ────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Launch — Examples 1 & 3 + Rules",
           framework_tag="teacher models", dok=2, minutes=12)
    card(s, "\U0001f4d8 nth Roots & Rational Exponents — Rules",
         [RULES_TEXT],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), color=GREEN)
    card(s, "Examples: 5-1-savvas-example-1-lesson-5-1  &  5-1-savvas-example-3-lesson-5-1",
         ["Two examples because this is a compressed lesson — both needed for bank/assessment coverage.",
          "Ex 1: evaluate nth roots and rational exponents.",
          "Ex 3: simplify and rationalize expressions with nth-root denominators."],
         Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), color=BLUE)

    # ── Explore ───────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Explore — Think-Pair-Share + ⭐ DOK-3 Capstone",
           framework_tag="student work", dok="2-3", minutes=33)
    text(s, "7 items in order. Plan ∼15 min for ⭐ Practice #50 DOK-3 Capstone.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Try It 1 · 5-1-savvas-try-it-1-lesson-5-1",
         "Evaluate nth roots and write in rational-exponent form.", False),
        ("Try It 3 · 5-1-savvas-try-it-3-lesson-5-1",
         "Simplify expression with rational exponent denominator. Rationalize.", False),
        ("Practice #28 · 5-1-savvas-q28",
         "Evaluate/simplify. Apply x^(1/n) = ⁿ√x.", False),
        ("Practice #30 · 5-1-savvas-q30",
         "Simplify rational exponent expression.", False),
        ("Practice #37 · 5-1-savvas-q37",
         "Rationalize an nth-root denominator.", False),
        ("Practice #48 · 5-1-savvas-q48",
         "Mixed: evaluate and simplify; connect back to even-n rule.", False),
        ("Practice #50  ⭐ DOK-3 CAPSTONE · 5-1-savvas-q50",
         "Cylinder volume V=πr²h with h=2r. Derive r from V, then find lateral surface area. "
         "Topic 5 Assessment Item 1A rehearsal. Plan ~15 min.", True),
    ]
    for lbl, txt, is_dok3 in items:
        color = ACCENT if is_dok3 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.78), color=color)
        top += Inches(0.84)

    # ── Share / Summary ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Share / Summary — DOK-3 Reveal",
           framework_tag="academic conversation", dok=2, minutes=5)
    card(s, "\U0001f4e2 Capstone Answer — 5-1-savvas-q50",
         ["r = (V/(2π))^(1/3).  Lateral SA = 4πr².",
          "",
          "For V = 169.65 ft³:  r ≈ 3 ft,  Lateral SA ≈ 113 ft².",
          "",
          "Sentence frame: \"I found r by ___, then substituted into ___ because ___.\""],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=BLUE)

    # ── Exit Ticket ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Exit Ticket — SAT/ACT Reinforce",
           framework_tag="not graded", dok="1-2", minutes=3)
    card(s, "✍️ Practice #49 · 5-1-savvas-q49 (SAT/ACT format)",
         ["Complete on your exit ticket.",
          "",
          "1. Today I learned that ___.",
          "2. The rule I used most was ___ because ___.",
          "3. One thing I want to review: ___."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


if __name__ == "__main__":
    p = build_p1("L51_P1_Slides.pptx")
    print(f"Built {p}")
