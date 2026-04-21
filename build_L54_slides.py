"""L54_P{1,2,3}_Slides.pptx — projection decks for Lesson 5-4 Klimsara-adapted.

Three decks, one per period. Same slide helpers as build_L43_slides.py.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, period_label, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, f"Algebra 2  ·  Lesson 5-4  ·  {period_label}",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
         size=28, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Klimsara-adapted · Student-centered · No Blooket",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


# ── P1 deck ───────────────────────────────────────────────────────────────

def build_p1(path="L54_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 5-4 · Period 1",
                "Solving Radical Equations: Foundation")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Do Now — Explore & Reason",
           framework_tag="exploration", dok=3, minutes=5)
    card(s, "5-4-savvas-model-discuss-lesson-5-4-launch · Explore & Reason: 3(a+1)²+2=11",
         ["Look at the equation on the board.",
          "",
          "A) Solve using your current method.",
          "B) Compare methods with your partner.",
          "C) What if the power were ½ or ⅓ instead of 2?"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Launch — Examples 1 & 3 + Rules",
           framework_tag="teacher models", dok=2, minutes=12)
    card(s, "\U0001f4d8 Solving Radical Equations — Rules",
         ["1. Isolate the radical.",
          "2. Raise both sides to the index power.",
          "3. Solve.",
          "4. CHECK in the ORIGINAL equation — reject extraneous."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), color=GREEN)
    card(s, "Examples: 5-4-savvas-example-1-lesson-5-4  &  5-4-savvas-example-3-lesson-5-4",
         ["Model Ex 1 fully (one-radical); release into Ex 3 (extraneous root).",
          "Emphasize: squaring is NOT reversible — always check in original."],
         Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), color=BLUE)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Explore — Think-Pair-Share",
           framework_tag="student work", dok="1-2", minutes=33)
    text(s, "6 items in order. Use the Radical Equation Rules card on your packet.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.1)
    items = [
        ("Try It 1 · 5-4-savvas-try-it-1-lesson-5-4", "Isolate radical, square both sides, check in original."),
        ("Try It 3 · 5-4-savvas-try-it-3-lesson-5-4", "Solve — watch for extraneous solution."),
        ("Practice #22 · 5-4-savvas-q22", "Solve the radical equation. Check both possible solutions."),
        ("Practice #8 · 5-4-savvas-q8", "Solve. x=9 valid; x=−1 extraneous — verify."),
        ("Practice #15 · 5-4-savvas-q15", "Identify any extraneous solutions."),
        ("Practice #4 · 5-4-savvas-q4", "Nazim's error: lost x=6, kept only −3. Find the mistake."),
    ]
    for lbl, txt in items:
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.65), color=BLUE)
        top += Inches(0.72)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Share / Summary",
           framework_tag="academic conversation", dok=2, minutes=5)
    card(s, "\U0001f4e2 Answer key highlights",
         ["q8: x=9 is valid; x=−1 is extraneous (fails check in original).",
          "q4: Nazim lost solution x=6 — only kept −3 which fails check.",
          "Cold-call q15: student explains extraneous-identification step-by-step."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.8), color=BLUE)
    card(s, "\U0001f52d Preview of Period 2",
         ["Next class: rational exponents — x^(m/n) = (n√x)^m.",
          "Same isolate-then-raise move, but with reciprocal exponent.",
          "Period 2 ends with ⭐ DOK-3 Half-Life capstone (Topic 5 spine)."],
         Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.3), color=ACCENT)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Exit Ticket — Summary Recap",
           framework_tag="not graded", dok="1-2", minutes=2)
    card(s, "✍️ Complete on your exit ticket",
         ["1. Today I learned that ___.",
          "2. An extraneous solution is ___ because ___.",
          "3. One thing I'm still unsure about is ___."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)
    text(s, "Reinforce: 5-4-savvas-q21 (∛x+8=13) as optional HW",
         Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4), size=13, color=GRAY)

    prs.save(path)
    return path


# ── P2 deck ───────────────────────────────────────────────────────────────

def build_p2(path="L54_P2_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 5-4 · Period 2",
                "Rational Exponents + ⭐ DOK-3 Half-Life (Topic 5 Spine)")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Do Now — Solve Cube Root",
           framework_tag="bridge from P1", dok=1, minutes=5)
    card(s, "Practice #21 · 5-4-savvas-q21",
         ["∛x + 8 = 13 — solve.",
          "",
          "Same isolate-then-raise move from P1, cube root instead of square root.",
          "Bridge: what reciprocal power undoes a cube root?"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Launch — Example 4 (Rational Exponents)",
           framework_tag="teacher models", dok=2, minutes=12)
    card(s, "\U0001f4d8 Rational Exponent Rules",
         ["1. x^(m/n) = (ⁿ√x)^m.",
          "2. Raise both sides to reciprocal power n/m.",
          "3. Check — even denominator → possible extraneous.",
          "4. Half-life: y = a·(0.5)^(t/h).",
          "",
          "Example (5-4-savvas-example-4-lesson-5-4): model steps on board.",
          "  Show rational exponent → raise to reciprocal → solve."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GREEN)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Explore — TPS + ⭐ DOK-3 Capstone",
           framework_tag="student work", dok="2-3", minutes=33)
    text(s, "5 items in order. Plan ∼12 min for Practice #41 ⭐ DOK-3 Spine (Half-Life).",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Try It 4 · 5-4-savvas-try-it-4-lesson-5-4", "Rational exponent equation. Raise to reciprocal, check."),
        ("Practice #12 · 5-4-savvas-q12", "Solve rational exponent equation. State any restrictions."),
        ("Practice #33 · 5-4-savvas-q33", "Solve. Check for extraneous (even denominator?)."),
        ("Practice #34 · 5-4-savvas-q34", "Solve. Full solution — identify reciprocal exponent."),
        ("Practice #41  ⭐ DOK-3 CAPSTONE · 5-4-savvas-q41",
         "Half-life: y = 50·(0.5)^(t/5). Find remaining amount after 16 hours. Answer ≈ 5.44 mL."),
        ("Reinforce · 5-4-savvas-q19", "Secondary DOK-3 — extraneous-argument written HW."),
    ]
    for i, (lbl, txt) in enumerate(items):
        color = ACCENT if i == 4 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.78), color=color)
        top += Inches(0.84)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Share / Summary",
           framework_tag="DOK-3 reveal + P3 bridge", dok=2, minutes=5)
    card(s, "\U0001f4e2 Answer key · DOK-3 Half-Life Capstone",
         ["q41 ≈ 5.44 mL.",
          "Interpretation: about 5.44 mL of the 50 mL drink remains after 16 hours.",
          "Work: 50·(0.5)^(16/5) = 50·(0.5)^3.2 ≈ 5.44"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), color=BLUE)
    card(s, "\U0001f52d Bridge to Period 3",
         ["Next class: solve for a variable (strip outer ops LIFO) + table-evaluation.",
          "Period 3 rehearses assessment Item 1A (solve for variable) and Item 2 (table).",
          "Know reciprocal exponent move cold."],
         Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), color=ACCENT)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Exit Ticket — Summary Recap",
           framework_tag="not graded", dok="1-2", minutes=3)
    card(s, "✍️ Complete on your exit ticket",
         ["1. Today I learned that ___.",
          "2. When the exponent denominator is even, I must ___ because ___.",
          "3. The half-life formula step I found hardest was ___."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)
    text(s, "Reinforce: 5-4-savvas-q19 extraneous-argument as written HW",
         Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4), size=13, color=GRAY)

    prs.save(path)
    return path


# ── P3 deck ───────────────────────────────────────────────────────────────

def build_p3(path="L54_P3_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 5-4 · Period 3",
                "Assessment-Critical: Solve for Variable + Modeling")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Do Now — Higher Order Thinking",
           framework_tag="bridge from P2", dok=3, minutes=5)
    card(s, "Practice #19 · 5-4-savvas-q19",
         ["Higher Order Thinking: when do rational-exponent equations",
          "have extraneous solutions?",
          "",
          "Answer: ONLY when the denominator of the exponent is EVEN.",
          "Bridge: today’s assessment items use cube root (odd denom) — no extraneous risk."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Launch — Ex 2 + Q44 (paired, assessment-critical)",
           framework_tag="teacher models 2 examples", dok=2, minutes=12)
    card(s, "\U0001f4d8 Ex 2 · 5-4-savvas-example-2-lesson-5-4 (Rewrite Formula — Golden Gate)",
         ["▶ Isolate a variable: strip outer operations LIFO.",
          "Teacher models domain-restriction notation step-by-step."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.2), color=GREEN)
    card(s, "⭐ Q44 · 5-4-savvas-q44 (Table Completion — Assessment-Rehearsal)",
         ["RULES for Table-Eval:",
          "  1. Substitute V value.",
          "  2. Compute, round to nearest tenth.",
          "  3. The table is a point-cloud of the graph.",
          "Teacher walks through first row, students complete remaining rows.",
          "⭐ IDENTICAL structure to assessment Item 2 answer key."],
         Inches(0.6), Inches(4.1), Inches(12.1), Inches(3.0), color=ACCENT)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Explore — TPS · Assessment Rehearsal",
           framework_tag="student work · pure DOK-2", dok=2, minutes=33)
    text(s, "6 items. Q10 and Q25–27 = Item 1A rehearsals. Q39 + Q44 = Item 2 rehearsals.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Try It 2 · 5-4-savvas-try-it-2-lesson-5-4", "Rewrite formula — isolate the radical variable."),
        ("Practice #10 · 5-4-savvas-q10", "Solve for variable. Assessment Item 1A rehearsal."),
        ("Practice #25 · 5-4-savvas-q25", "Solve for variable. State domain / restrictions."),
        ("Practice #26 · 5-4-savvas-q26", "Solve for variable. Rationalize pattern."),
        ("Practice #27 · 5-4-savvas-q27", "Solve for variable. Full algebraic justification."),
        ("Practice #39 · 5-4-savvas-q39", "Substitute + evaluate. Item 2 table rehearsal."),
        ("Reinforce · 5-4-savvas-q46  ⭐ ASSESSMENT-REHEARSAL",
         "Escape Velocity performance task — Topic 5 LEHS rehearsal."),
    ]
    for i, (lbl, txt) in enumerate(items):
        color = ACCENT if i == 6 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.65), color=color)
        top += Inches(0.72)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Share / Summary · Concept Summary",
           framework_tag="assessment preview", dok=2, minutes=5)
    card(s, "\U0001f4ca Q44 Table Answer Key",
         ["V = {0, 0.25, 2, 16, 32, 64, 128}",
          "x = {0, 0.5, 1, 2, 2.5, 3.2, 4}",
          "",
          "IDENTICAL to assessment Item 2 answer key.",
          "\U0001f3af Topic 5 assessment: expect Item 1A on solve-for-variable + Item 2 on table evaluation."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=BLUE)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Exit Ticket — Pre-Assessment Confidence",
           framework_tag="confidence check", dok="1-2", minutes=3)
    card(s, "✍️ Rate yourself 1–5 + name one review item",
         ["• I can solve a radical equation and check for extraneous roots.  Confidence: 1 / 2 / 3 / 4 / 5",
          "• I can solve for a variable using reciprocal exponents.  Confidence: 1 / 2 / 3 / 4 / 5",
          "",
          "One thing I want to review before Topic 5 assessment: ___"],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)
    text(s, "Reinforce: 5-4-savvas-q46 (Escape Velocity performance task)",
         Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4), size=13, color=GRAY)

    prs.save(path)
    return path


if __name__ == "__main__":
    for fn in [build_p1, build_p2, build_p3]:
        p = fn()
        print(f"Built {p}")
