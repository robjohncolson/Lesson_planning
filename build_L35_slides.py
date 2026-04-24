"""L35_P2_Slides.pptx + L35_P3_Slides.pptx — projection decks for 
Lesson 3-5 Periods 2 and 3. Phase slides show the same items the student packet
contains, so students scanning the screen see what they're working on.

No Blooket phase. Single builder, two decks.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


# ── helpers ───────────────────────────────────────────────────────────────

def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, period_label, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, f"Algebra 2  ·  Lesson 3-5  ·  {period_label}",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    # Right-side badges
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


# ── slide builders ────────────────────────────────────────────────────────

def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
         size=28, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Student-centered · No Blooket",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


def do_now_slide(prs, period_label, qid, bridge=None):
    q = qb.get(qid)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, period_label, "Do Now — 5 min", framework_tag="bridge", dok=1, minutes=5)
    text(s, q["prompt"], Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.5),
         size=22, color=DARK)
    if q.get("answers"):
        lines = [f"({chr(65+i)})  {a}" for i, a in enumerate(q["answers"])]
        text(s, lines, Inches(1.0), Inches(4.3), Inches(11), Inches(2),
             size=20, color=DARK)
    if bridge:
        card(s, "🔗 Bridge from yesterday", bridge,
             Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.4), color=ACCENT)


def launch_slide(prs, period_label, launch_title, rule_title, rule_lines, *,
                 minutes=12):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, period_label, f"Launch — {launch_title}", framework_tag="teacher models", dok=2, minutes=minutes)
    card(s, rule_title, rule_lines,
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GREEN)


def explore_slide(prs, period_label, items_list, *, minutes=33):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, period_label, "Explore — Think-Pair-Share",
           framework_tag="student work", dok=2, minutes=minutes)
    text(s, "Work with a partner. Use the rule card from Launch.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=18, color=GRAY)
    top = Inches(2.2)
    for label, prompt in items_list:
        card(s, label, [prompt], Inches(0.6), top, Inches(12.1), Inches(1.0),
             color=BLUE)
        top += Inches(1.1)


def share_slide(prs, period_label, frame_text, bridge=None, *, minutes=5):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, period_label, "Share / Summary", framework_tag="academic conversation", dok=2, minutes=minutes)
    card(s, "📢 Sentence frame (one pair reads aloud)", frame_text,
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.5), color=BLUE)
    if bridge:
        card(s, "🔭 Bridge to tomorrow", bridge,
             Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.6), color=ACCENT)


def exit_slide(prs, period_label, stems, *, minutes=2):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, period_label, "Exit Ticket — Summary Recap", framework_tag="not graded", dok="1-2", minutes=minutes)
    card(s, "✍️ Complete the stems on your exit ticket", stems,
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)


# ── Deck builders ────────────────────────────────────────────────────────

def build_p2(path="L35_P2_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 3-5 · Period 2",
                "Multiplicity and Factored-Form Graphing")

    do_now_slide(prs, "Period 2 (Tue)", "3-5-savvas-q13")

    launch_slide(prs, "Period 2 (Tue)",
                 "Multiplicity & the Graph",
                 "📘 Multiplicity Rules",
                 [
                     "• Odd multiplicity → graph CROSSES the x-axis at that zero",
                     "• Even multiplicity → graph is TANGENT to the x-axis (touches & turns back)",
                     "",
                     "Example: f(x) = (x − 2)³(x + 1)²",
                     "  x = 2 has multiplicity 3 (odd) → CROSSES",
                     "  x = −1 has multiplicity 2 (even) → TANGENT",
                 ])

    # Explore: show the 5 Practice items at a glance
    items = [
        ("Try It 2a", qb.get("3-5-tryit-2a")["prompt"][:140]),
        ("Try It 2b", qb.get("3-5-tryit-2b")["prompt"][:140]),
        ("Practice #14", qb.get("3-5-savvas-q14")["prompt"][:140]),
        ("Practice #15", qb.get("3-5-savvas-q15")["prompt"][:140]),
        ("Practice #16", qb.get("3-5-savvas-q16")["prompt"][:140]),
    ]
    explore_slide(prs, "Period 2 (Tue)", items, minutes=33)

    share_slide(prs, "Period 2 (Tue)",
                ["\"For f(x) = ___, the zero x = ___ has multiplicity ___,",
                 " so the graph ___ the x-axis at that zero.\""],
                bridge=[
                    "If a cubic has zeros 2, 1 + √3, and 1 − √3,",
                    "what is the nature of its coefficients —",
                    "rational, irrational, or complex? Why?",
                    "",
                    "(Direct prep for Topic 3 Assessment Q#2)",
                ])

    exit_slide(prs, "Period 2 (Tue)",
               [
                   "1. Today I learned that ___.",
                   "2. The rule I want to remember is ___.",
                   "3. One thing I'm still unsure about is ___.",
               ])

    prs.save(path)
    return path


def build_p3(path="L35_P3_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 3-5 · Period 3",
                "Real + Complex Zeros, and Modeling")

    do_now_slide(prs, "Period 3 (Wed)", "3-5-savvas-q11",
                 bridge=[
                     "A cubic has zeros 2, 1 + √3, and 1 − √3.",
                     "What's the nature of its coefficients —",
                     "rational, irrational, or complex? Why?",
                 ])

    launch_slide(prs, "Period 3 (Wed)",
                 "Complex Zeros + Modeling",
                 "📘 Complex-Zero Conjugate Rule + Modeling",
                 [
                     "• A polynomial with REAL coefficients has complex zeros in CONJUGATE PAIRS.",
                     "• If (a + bi) is a zero, then (a − bi) must also be a zero.",
                     "• A degree-n polynomial has exactly n zeros (with multiplicity).",
                     "",
                     "📐 Modeling (Storage Box):",
                     "• Factored form V(x) = (expr₁)(expr₂)(expr₃) — each factor is a dimension.",
                     "• Use physical constraint (\"length is greatest\") to assign dimensions.",
                     "• Solve cubic equation V(x) = 10 ft³ — show all work.",
                 ])

    items = [
        ("Try It 3a", qb.get("3-5-tryit-3a")["prompt"][:140]),
        ("Try It 3b", qb.get("3-5-tryit-3b")["prompt"][:140]),
        ("Practice #17", qb.get("3-5-savvas-q17")["prompt"][:140]),
        ("Practice #27  ⭐ DOK-3 DRIVER", qb.get("3-5-savvas-q27")["prompt"][:180]),
    ]
    explore_slide(prs, "Period 3 (Wed)", items, minutes=38)

    share_slide(prs, "Period 3 (Wed)",
                ["\"For the Storage Box, the factored form was ___.",
                 " I chose ___ as the length because ___.",
                 " To find V = 10 ft³, I ___.\""],
                bridge=[
                    "Tomorrow's assessment Q#6 (Lucy's tray)",
                    "uses the same factored-volume idea.",
                    "Different physical setup → same strategy.",
                ])

    exit_slide(prs, "Period 3 (Wed)",
               [
                   "1. Today I learned that ___.",
                   "2. For the Storage Box, the hardest step was ___.",
                   "3. Thursday's assessment will test ___.",
               ])

    prs.save(path)
    return path


if __name__ == "__main__":
    p2 = build_p2()
    p3 = build_p3()
    print(f"Built {p2}")
    print(f"Built {p3}")
