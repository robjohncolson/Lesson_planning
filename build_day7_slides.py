"""Day_7_Slides.pptx — Combined Day 7 projection deck (Savvas-only).

Eleven slides: title + 9 phase slides + wrap.
DOK-3 driver = Savvas Practice #30 (Venetta deli profit), Savvas-declared
performance-task anchor for this lesson.
"""
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY  = RGBColor(0x0B, 0x3C, 0x6C)
BLUE  = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
DARK  = RGBColor(0x11, 0x22, 0x33)
GRAY  = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
TEAL  = RGBColor(0x00, 0x97, 0xA7)


def add_background(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_badge(slide, text, left, top, *, color=NAVY,
              text_color=RGBColor(0xFF, 0xFF, 0xFF)):
    width = Inches(0.28 + 0.11 * len(text))
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, Inches(0.35))
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shp, width


def add_header(slide, phase_num, total, phase_title, framework, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    add_text(slide, phase_title, Inches(0.5), Inches(0.18),
             Inches(9.5), Inches(0.9),
             size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    meta = f"Phase {phase_num}/{total}  \u00b7  Algebra 2  \u00b7  Lesson 3-5  \u00b7  Day 7  (Synthesis)"
    add_text(slide, meta, Inches(0.5), Inches(0.82),
             Inches(9), Inches(0.35),
             size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    bx = Inches(10.3)
    by = Inches(0.4)
    if framework:
        shp, w = add_badge(slide, framework, bx, by,
                           color=RGBColor(0x22, 0x33, 0x44))
        bx = bx + w + Inches(0.08)
    if dok and dok != "\u2014":
        add_badge(slide, dok, bx, by, color=BLUE)

    if minutes is not None:
        add_text(slide, f"{minutes} min", Inches(10.3), Inches(0.85),
                 Inches(2.5), Inches(0.35),
                 size=14, bold=True,
                 color=RGBColor(0xCC, 0xDD, 0xEE))


def add_phase_slide(prs, phase_num, total, phase_title, framework, dok,
                    minutes, body_lines, footer=None, body_size=22):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, phase_num, total, phase_title, framework, dok, minutes)
    add_text(slide, body_lines,
             Inches(0.6), Inches(1.6),
             Inches(12.1), Inches(5.3),
             size=body_size, color=DARK)
    if footer:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     0, Inches(6.85), SLIDE_W, Inches(0.65))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        add_text(slide, footer, Inches(0.5), Inches(6.92),
                 Inches(12.3), Inches(0.5),
                 size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    return slide


def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---- Title slide ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s, NAVY)
    add_text(s, "ALGEBRA 2  \u00b7  LESSON 3-5",
             Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
             size=22, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(s, "Day 7  \u2014  Synthesis + Venetta Performance Task",
             Inches(0.6), Inches(1.75), Inches(12), Inches(1.4),
             size=34, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, [
        "Essential Question:",
        "How do zeros, sign charts, and the factored form work together",
        "to answer when a business is profitable, breaks even, or loses money?",
        "",
        "DOK-3 driver: Savvas Practice #30 (Venetta deli profit).",
        "All items Savvas-anchored.",
    ],
        Inches(0.6), Inches(3.4), Inches(12), Inches(3.0),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_text(s, "57-min lesson  \u00b7  Savvas-declared DOK 3",
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
             size=16, color=RGBColor(0x88, 0xAA, 0xCC))

    total = 9

    # ---- Phase 1: Do Now A ----
    add_phase_slide(prs, 1, total,
        "Do Now  \u2014  SAT-Style Graph ID",
        "Do Now A", "DOK 2", 5, [
            "SAVVAS PRACTICE #29  \u2014  SAT/ACT Style",
            "",
            "Without a graphing calculator, which graph shows",
            "          f(x) = x\u00b3 + x\u00b2 \u2212 4x  ?",
            "",
            "Circle:    A          B          C          D",
            "",
            "Write ONE sentence: how do you know?",
            "",
            "Hint: zeros \u2192 behavior \u2192 sketch",
            "",
            "Silent  \u00b7  Pencil only  \u00b7  No graphing calculator",
        ],
        footer="When everyone\u2019s in:  Blooket code on the next slide.")

    # ---- Phase 2: Blooket Login ----
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s2)
    add_header(s2, 2, total, "Blooket  \u2014  Log In", "Do Now B", "\u2014", 2)
    add_text(s2, "Blooket code:",
             Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
             size=28, color=GRAY)
    box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(1.5), Inches(2.6), Inches(10.3), Inches(3))
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = NAVY
    box.line.width = Pt(3)
    add_text(s2, "( write the code here )",
             Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.5),
             size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s2, "Chromebooks open  \u00b7  Type code  \u00b7  Nickname",
             Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             size=16, color=DARK, align=PP_ALIGN.CENTER)
    bar2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.85),
                               SLIDE_W, Inches(0.65))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT
    bar2.line.fill.background()
    add_text(s2, "2 minutes to log in.",
             Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.5),
             size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    # ---- Phase 3: Blooket Rule Recall ----
    add_phase_slide(prs, 3, total,
        "Blooket  \u2014  Synthesis Rule Recall",
        "Do Now C", "DOK 1", 6, [
            "Whole-lesson synthesis rules:",
            "",
            "  \u2022  Zero \u21d4 factor theorem: zero at c  \u21d4  (x \u2212 c)",
            "  \u2022  Odd multiplicity  \u2192  graph CROSSES the x-axis",
            "  \u2022  Even multiplicity  \u2192  graph is TANGENT to x-axis (turning point)",
            "  \u2022  Non-real complex zeros come in CONJUGATE PAIRS",
            "  \u2022  Sign chart: find zeros, pick test values, determine sign",
            "  \u2022  Polynomial inequality: solution = intervals where sign matches",
            "",
            "Play fast.  Teacher is watching the dashboard.",
        ],
        footer="Close game at 0.  Blooket Wrap next (1 min).")

    # ---- Phase 4: Blooket Wrap ----
    add_phase_slide(prs, 4, total,
        "Blooket Wrap  \u2014  Candy + Lowest Rule",
        "Do Now D", "DOK 1", 1, [
            "1.  Candy pass (~20 sec)",
            "      Walk the candy jar once.  Top finisher(s) get an extra piece.",
            "",
            "2.  Lowest-scored rule aloud (~25 sec)",
            "      Announce the lowest-percent rule from the dashboard.",
            "      Whole class repeats it in unison.  One volunteer again.",
            "",
            "3.  Transition (~15 sec)",
            "      \u201cPull out your Do Now sheet.  Launch coming up.\u201d",
        ],
        footer="\u201cPull out your Do Now.  8 minutes on Launch.\u201d")

    # ---- Phase 5: Launch ----
    add_phase_slide(prs, 5, total,
        "Launch  \u2014  3-Concept Rapid Review",
        "Launch", "DOK 2", 8, [
            "Three quick checks.  Call out answers.",
            "",
            "1.  MULTIPLICITY:  f(x) = (x \u2212 1)\u00b2(x + 3)",
            "      At x = 1:   even mult \u2192 TANGENT to x-axis (turning point)",
            "      At x = \u22123:  odd mult \u2192 CROSSES the x-axis",
            "",
            "2.  COMPLEX ZEROS:  f(x) = x\u00b3 \u2212 x\u00b2 \u2212 4x \u2212 6  [Savvas Practice #17]",
            "      Real zero x = 3; synthetic division \u2192 x\u00b2 + 2x + 2",
            "      Complex zeros:  x = \u22121 + i   and   x = \u22121 \u2212 i",
            "",
            "3.  INEQUALITY:  2x\u00b3 + 12x\u00b2 + 12x < 0  [Savvas Try It 6a]",
            "      Factor, find zeros, build sign chart.",
            "      Solution:  x < \u22123 \u2212 \u221a3  or  \u22123 + \u221a3 < x < 0",
        ],
        footer="\u201cPackets to the Venetta page.  22 minutes.  Read the rules on the page.\u201d",
        body_size=20)

    # ---- Phase 6: Venetta Explore (DOK 3) ----
    add_phase_slide(prs, 6, total,
        "Venetta Deli  \u2014  the DOK-3 driver",
        "Explore", "DOK 3", 22, [
            "SAVVAS PRACTICE #30  \u2014  Performance Task",
            "",
            "P(t) = t\u00b3 + t\u00b2 \u2212 6t     (profit in hundreds of $, t years since 2000)",
            "",
            "  Part A:  Sketch a graph.  Factor P(t) first.",
            "  Part B:  During which year(s) did franchises break even?  (P = 0)",
            "  Part C:  Interpret years of negative P.  Use the sign chart.",
            "  Part D:  During which year range was the franchise profitable?  (P > 0)",
            "",
            "All rules are on your packet page.  Teacher circulates with prompts only.",
        ],
        footer="Stuck after 10 min?  Hint card:  \u201cFactor first.  Which zeros have t \u2265 0?\u201d",
        body_size=22)

    # ---- Phase 7: Share / Summary ----
    add_phase_slide(prs, 7, total,
        "Share / Summary",
        "Share/Summary", "DOK 2", 4, [
            "Cold-call Pair 1:  describe the sign chart for P(t).",
            "",
            "Cold-call Pair 2:  justify why P(t) > 0 for t > 2.",
            "",
            "Self-rating on your packet:   \u2713  /  partly  /  not yet",
            "",
            "Preview Day 8:  quiz + differentiated stations.",
        ],
        footer="\u201cPencils up.  Exit ticket.  5 minutes.\u201d")

    # ---- Phase 8: Exit Ticket ----
    add_phase_slide(prs, 8, total,
        "Exit Ticket  \u2014  Synthesis Recap",
        "Exit Ticket", "DOK 1\u20132", 5, [
            "Fill in each blank:",
            "",
            "1.  Odd multiplicity \u2192 graph ___________ the x-axis.",
            "    Even multiplicity \u2192 graph is ___________ to x-axis.",
            "",
            "2.  Non-real complex zeros always come in ____________ pairs.",
            "",
            "3.  A sign chart tells you where a polynomial is _____ (> 0) or _____ (< 0).",
            "",
            "4.  Venetta\u2019s franchises were profitable for years _____ and beyond",
            "    because ____________________________________________________.",
            "",
            "Biggest thing I learned today across the whole lesson: _______________",
        ],
        footer="Packets to the front.  Pack up.  Day 8: quiz.",
        body_size=21)

    # ---- Phase 9: Blooket code placeholder (if needed for Day 7 game) ----
    # (Wrap slide — end of lesson)
    sw = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(sw, NAVY)
    add_text(sw, "Great work today.",
             Inches(0.6), Inches(2.1), Inches(12), Inches(1.2),
             size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(sw, [
        "Collect:  Do Now  \u00b7  Venetta packet  \u00b7  Exit Ticket",
        "",
        "Day 8:  quiz + differentiated stations.",
        "Review tonight:  zeros, multiplicity, sign charts, Venetta.",
    ],
        Inches(0.6), Inches(3.8), Inches(12), Inches(2.4),
        size=22, color=RGBColor(0xDD, 0xEE, 0xFF),
        align=PP_ALIGN.CENTER)

    prs.save(path)


if __name__ == "__main__":
    build("Day_7_Slides.pptx")
    print("Built Day_7_Slides.pptx")
