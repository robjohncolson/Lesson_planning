"""L44_P1_Slides.pptx — projection deck for Lesson 4-4 (single-period quick).

Adding & Subtracting Rational Expressions + DOK-3 Electrical Resistance capstone
(Topic 4 LEHS Q#15 prep). Compressed 2-example Launch (Ex 3 add + Ex 4 subtract).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


# ── ID manifest (mirrors build_L44_P1_packets.py) ──────────────────────────

DO_NOW_ID   = "4-4-savvas-model-discuss-lesson-4-4-launch"
LAUNCH_IDS  = [
    "4-4-savvas-example-3-lesson-4-4",
    "4-4-savvas-example-4-lesson-4-4",
]
EXPLORE_IDS = [
    "4-4-savvas-try-it-3-lesson-4-4",
    "4-4-savvas-try-it-4-lesson-4-4",
    "4-4-savvas-q12",
    "4-4-savvas-q16",
    "4-4-savvas-q26",
    "4-4-savvas-q32",
]
REINFORCE_IDS = ["4-4-savvas-q34"]

_ALL_IDS = [DO_NOW_ID] + LAUNCH_IDS + EXPLORE_IDS + REINFORCE_IDS
qb.get_for_packet(_ALL_IDS)


# ── Drawing primitives (cloned from build_L51_slides.py) ───────────────────

def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, "Algebra 2  ·  Lesson 4-4  ·  Single Period (Quick)",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.8),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.85), Inches(11.7), Inches(0.8),
         size=24, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Student-centered · Topic 4 LEHS Q#15 rehearsal",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


RULES_CARD_LINES = [
    "1. LIKE denominators: combine numerators, keep the denominator, simplify.",
    "2. UNLIKE: factor each denom → find LCM → rewrite each fraction over LCM → combine numerators → simplify.",
    "3. LCM takes every DISTINCT factor at its HIGHEST power across both denominators.",
    "4. After combining, always check for a common factor to cancel.",
]


# ── P1 deck ────────────────────────────────────────────────────────────────

def build_p1(path="L44_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(
        prs,
        "Lesson 4-4 · Quick — Adding & Subtracting Rational Expressions",
        "Add/Subtract Rational Expressions + ⭐ Electrical Resistance DOK-3 (Q#15 Prep)",
    )

    # ── Do Now ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Do Now — Model & Discuss: 1/2 + 1/3 conflict",
           framework_tag="exploration", dok=1, minutes=5)
    card(s, f"{DO_NOW_ID}",
         ["Two students add 1/2 + 1/3 in conflicting ways. Which is right, and why?",
          "",
          "A) What's the LCM of 2 and 3?",
          "B) Why can't we just add numerators and denominators?",
          "C) Bridge: yesterday we MULTIPLIED rational expressions. What NEW move will ADDING require that multiplying didn't?"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    # ── Launch ────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Launch — Examples 3 & 4 + Rules [COMPRESSED]",
           framework_tag="teacher models", dok=2, minutes=13)
    card(s, "\U0001f4d8 Add/Subtract Rational Expressions — Rules",
         RULES_CARD_LINES,
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), color=GREEN)
    card(s, "Examples: 4-4-savvas-example-3-lesson-4-4  &  4-4-savvas-example-4-lesson-4-4",
         ["Ex 3: ADD with unlike denominators — factor, build LCM, rewrite, combine.",
          "Ex 4: SUBTRACT with unlike denominators — SAME setup, but DISTRIBUTE the minus sign across EVERY term of the second numerator.",
          "KEY TRAP: sign distribution on subtraction. #1 error source — watch it on Ex 4."],
         Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.6), color=BLUE)

    # ── Explore ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Explore — Think-Pair-Share + ⭐ DOK-3 Resistance",
           framework_tag="student work", dok="2-3", minutes=32)
    text(s, "6 items in order. Plan ~10 min for ⭐ Practice #32 (Electrical Resistance DOK-3).",
         Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.1)
    items = [
        ("Try It 3 · 4-4-savvas-try-it-3-lesson-4-4",
         "Add rational expressions (unlike denominators). Build the LCM first.", BLUE),
        ("Try It 4 · 4-4-savvas-try-it-4-lesson-4-4",
         "Subtract rational expressions (unlike denominators). Watch the sign distribution.", BLUE),
        ("Practice #12 · 4-4-savvas-q12",
         "LIKE denominators — quick confidence build.", BLUE),
        ("Practice #16 · 4-4-savvas-q16",
         "Unlike denominators — factor first, then LCM.", BLUE),
        ("Practice #26 · 4-4-savvas-q26",
         "Subtract with SIGN-DISTRIBUTION TRAP. Distribute the minus across EVERY term.", BLUE),
        ("Practice #32 ⭐ DOK-3 ELECTRICAL RESISTANCE · 4-4-savvas-q32",
         "Parallel circuits: 1/R_total = 1/R₁ + 1/R₂. Form the ratio R_A/R_B as a single simplified rational expression. Topic 4 LEHS Q#15 rehearsal. Plan ~10 min.", ACCENT),
    ]
    row_h = Inches(0.72)
    for lbl, txt, clr in items:
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), row_h, color=clr)
        top += row_h + Inches(0.05)

    # ── DOK-3 Spotlight ───────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "⭐ DOK-3 SPOTLIGHT — Electrical Resistance (Practice #32)",
           framework_tag="DOK-3 driver", dok=3, minutes=10)
    card(s, "\U0001f50c PARALLEL-RESISTOR FORMULA",
         ["For two resistors in parallel:   1/R_total = 1/R₁ + 1/R₂",
          "",
          "This is a SUM of rational expressions. Same LCM move as Ex 3.",
          "",
          "TASK: For circuit A and circuit B (each with two parallel resistors), form the RATIO R_A/R_B as a SINGLE simplified rational expression."],
         Inches(0.5), Inches(1.55), Inches(6.1), Inches(5.7), color=ACCENT)
    card(s, "\U0001f4ac CER EXPECTATION",
         ["(C)  Ratio as a single simplified rational expression.",
          "(E)  Algebra steps — common denom, add, simplify.",
          "(R)  Parallel-resistance formula + observation that the ratio depends only on the constituent resistances.",
          "",
          "STUCK? Ask: \"What's the LCM of the two denominators?\" Once they have that, the rest is algebra.",
          "",
          "BRIDGE: same structural reasoning as Topic 4 LEHS Q#15 (operations / closure)."],
         Inches(6.7), Inches(1.55), Inches(6.1), Inches(5.7), color=GOLD)

    # ── Share / Summary ───────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Share / Summary — Resistance Reveal + Q#15 Bridge",
           framework_tag="academic conversation", dok=2, minutes=5)
    card(s, "\U0001f4e2 Answer Key — Practice #32 (DOK-3 Resistance)",
         ["For 1/R_total = 1/R₁ + 1/R₂:",
          "  Common denominator R₁R₂  ⇒  1/R_total = (R₂ + R₁)/(R₁R₂)",
          "  Flip:  R_total = R₁R₂ / (R₁ + R₂)",
          "",
          "Form the ratio R_A/R_B as a SINGLE simplified rational expression using that form for each circuit.",
          "",
          "Sentence frame: \"For the parallel circuit, 1/R_total = 1/R₁ + 1/R₂. I found a common denominator of ___, added the reciprocals to get ___, then flipped to get R_total = ___.\"",
          "",
          "\U0001f517 Topic 4 LEHS Q#15 bridge: same LCM + operations/closure move you just did."],
         Inches(0.6), Inches(1.55), Inches(12.1), Inches(5.7), color=BLUE)

    # ── Exit Ticket ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Exit Ticket — Summary Recap",
           framework_tag="not graded", dok="1-2", minutes=3)
    card(s, "✍️ Summary Exit — 3 stems",
         ["1. To add rational expressions with unlike denominators, I must first find the ___ of the denominators.",
          "",
          "2. Practice #32 (resistance) showed me that reciprocal-sum problems (1/R₁ + 1/R₂) use the same algebraic steps as adding rational expressions because ___.",
          "",
          "3. One biggest thing I learned today: ___________________________"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    prs.save(path)
    print(f"Built {path}")
    return path


if __name__ == "__main__":
    build_p1("L44_P1_Slides.pptx")
