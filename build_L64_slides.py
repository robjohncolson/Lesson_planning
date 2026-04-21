"""L64_P1_Slides.pptx — projection deck for Lesson 6-4 (single-period).

Single deck: Logarithmic Functions + DOK-3 Sales Revenue Inverse (Item 28).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, period_label, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, f"Algebra 2  ·  Lesson 6-4  ·  {period_label}",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
         size=28, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Klimsara-adapted · Student-centered · No Blooket",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


RULES_TEXT = (
    "log_b 1 = 0.  log_b b = 1.  "
    "y = log_b x: Domain x>0, VA x=0, x-int (1,0), passes (b,1).  "
    "y = log_b(x−h)+k: VA x=h, Domain x>h.  "
    "Inverse: y=a·b^x ⇒ y=log_b(x/a);  "
    "y=log_b(x−h)+k ⇒ y=b^(x−k)+h.  "
    "Composition check: f(f⁻¹(x))=x."
)


# ── P1 deck (single period) ───────────────────────────────────────────────────

def build_p1(path="L64_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(
        prs,
        "Lesson 6-4 · Quick — Logarithmic Functions + ⭐ Sales Revenue DOK-3 (Item 28)",
        "Graphing Logs + Key Features + Inverses of Exp/Log Functions",
    )

    # ── Do Now ────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Do Now — 3 Items",
           framework_tag="bridge", dok="1-2", minutes=7)
    card(s, "6-4-savvas-q31 · Translation identification: h(x)=ln(x+2)−1",
         ["Multi-select: which transformations apply to h(x) = ln(x+2) − 1?",
          "",
          "Key: +2 inside → left shift 2. −1 outside → down 1.",
          "VA moves from x=0 to x=−2."],
         Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.5), color=GOLD)
    card(s, "6-4-savvas-q5 · Graph y=log₄ x; state key features",
         ["Graph on coordinate plane.",
          "State: domain, range, VA, x-intercept, end behavior as x→0⁺ and x→∞."],
         Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.8), color=BLUE)
    card(s, "6-4-savvas-q32 · SAT/ACT: inverse of f(x)=5^(x+1)",
         ["Find f⁻¹(x).",
          "",
          "Strategy: swap x↔y → x=5^(y+1) → log₅ x = y+1 → y = log₅ x − 1."],
         Inches(6.7), Inches(1.7), Inches(6.0), Inches(2.5), color=ACCENT)

    # ── Launch ────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Launch — Ex 3B: Inverse of a Log + Composition Check",
           framework_tag="teacher models", dok=2, minutes=8)
    card(s, "\U0001f4d8 Log Key Features + Inverse Rules",
         [RULES_TEXT],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), color=GREEN)
    card(s, "Example 3B: 6-4-savvas-example-3-lesson-6-4",
         ["g(x) = log₇(x+5). Find g⁻¹(x) and verify with composition.",
          "Step 1: y = log₇(x+5). Step 2: swap → x = log₇(y+5).",
          "Step 3: 7^x = y+5. Step 4: g⁻¹(x) = 7^x − 5.",
          "Check: g(g⁻¹(x)) = log₇(7^x−5+5) = log₇(7^x) = x. ✔",
          "90-sec contrast: Ex 3A (inverting exponential) → swap → log form."],
         Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.6), color=BLUE)

    # ── Explore ───────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Explore — Think-Pair-Share + ⭐ DOK-3 Sales Revenue",
           framework_tag="student work", dok="2-3", minutes=25)
    text(s, "5 items in order. Start with ⭐ q28 DOK-3 (∼10 min). 45-min Wed F: skip q9.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Practice #28 · 6-4-savvas-q28",
         "Sales Revenue R = 12·log(a+1)+25. Find inverse: a = 10^((R−25)/12) − 1. "
         "Judge which form is easier for specific R values. Plan ∼10 min.", True),
        ("Practice #21 · 6-4-savvas-q21",
         "Inverse of 5^(x−3). [Form B Q7 rehearsal]", False),
        ("Practice #24 · 6-4-savvas-q24",
         "Inverse of log₂(8x), token drag. [Form B Q12 rehearsal]", False),
        ("Practice #13 · 6-4-savvas-q13",
         "Asymptote of translated log from graph. [Form B Q11 rehearsal]", False),
        ("Practice #9 · 6-4-savvas-q9",
         "Are graphs inverses? Inverse from graph. [Form B Q9/Q10 rehearsal] "
         "⚠️ 45-min Wed F: skip this item.", False),
    ]
    for lbl, txt, is_dok3 in items:
        color = ACCENT if is_dok3 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.78), color=color)
        top += Inches(0.84)

    # ── Share / Summary ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Share / Summary — DOK-3 Reveal + 6-5 Bridge",
           framework_tag="academic conversation", dok=2, minutes=5)
    card(s, "\U0001f4e2 Capstone Answer — 6-4-savvas-q28",
         ["R = 12·log(a+1)+25  ⟹  a = 10^((R−25)/12) − 1.",
          "",
          "For R=49: (49−25)/12 = 2  →  a = 10²−1 = 99.",
          "For R=61: (61−25)/12 = 3  →  a = 10³−1 = 999.",
          "",
          "Inverse form easier for computation when R is given. Original easier for checking.",
          "",
          "Sentence frame: \"I isolated a by ___, giving a = ___. "
          "The inverse form is easier when ___ because ___.\"",
          "",
          "Bridge → 6-5: log properties (product/quotient/power) simplify these steps."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=BLUE)

    # ── Exit Ticket ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Exit Ticket — Inverse + Asymptotes",
           framework_tag="not graded", dok="1-2", minutes=5)
    card(s, "✍️ Custom Exit — Two-Part",
         ["1. Find the inverse of f(x) = 3·log₅(x − 2) + 1.",
          "",
          "2. State the vertical asymptote of f.",
          "",
          "3. State the horizontal asymptote of f⁻¹.",
          "",
          "4. One biggest thing I learned today: ___.",
          "",
          "ANSWER (TE): f⁻¹(x) = 5^((x−1)/3) + 2.  VA of f: x=2.  HA of f⁻¹: y=2."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


if __name__ == "__main__":
    p = build_p1("L64_P1_Slides.pptx")
    print(f"Built {p}")
