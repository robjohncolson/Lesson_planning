"""L43_P{1,2,3}_Slides.pptx — projection decks for Lesson 4-3 .

Three decks, one per period. Same slide helpers as build_L41_slides.py.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, period_label, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, f"Algebra 2  ·  Lesson 4-3  ·  {period_label}",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
         size=28, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Student-centered",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


# ── P1 deck ───────────────────────────────────────────────────────────────

def build_p1(path="L43_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 4-3 · Period 1",
                "Rational Expressions: Equivalent & Simplify")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Do Now — Notice & Wonder",
           framework_tag="exploration", dok=2, minutes=5)
    card(s, "4-3-savvas-model-discuss-lesson-4-3-launch · Notice & Wonder",
         ["Look at the expression on the board.",
          "",
          "A) Notice: 1 thing you notice",
          "B) Wonder: 1 thing you wonder",
          "C) Predict: What values of x cause a problem?"],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Launch — Examples 1 & 2 + Rules",
           framework_tag="teacher models", dok=2, minutes=12)
    card(s, "📘 Simplifying Rational Expressions — Rules",
         ["1. Factor numerator AND denominator fully.",
          "2. Cancel FACTORS, never terms across + or −.",
          "3. Domain excludes every zero of the ORIGINAL denominator",
          "   — including canceled factors."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), color=GREEN)
    card(s, "Examples: 4-3-savvas-example-1-lesson-4-3  &  4-3-savvas-example-2-lesson-4-3",
         ["Model each step: factor → identify restrictions → cancel → simplify.",
          "State the domain restriction(s) for each example."],
         Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), color=BLUE)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Explore — Think-Pair-Share",
           framework_tag="student work", dok="1-2", minutes=33)
    text(s, "10 items in order. Use the Simplifying Rules card on your packet.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.1)
    items = [
        ("Try It 1 · 4-3-savvas-try-it-1-lesson-4-3", "Simplify the rational expression. State domain restrictions."),
        ("Try It 2 · 4-3-savvas-try-it-2-lesson-4-3", "Simplify the rational expression. State domain restrictions."),
        ("Practice #20 · 4-3-savvas-q20", "Simplify. Identify canceled factor; state excluded value."),
        ("Practice #22 · 4-3-savvas-q22", "Simplify. Factor both numerator and denominator first."),
        ("Practice #24 · 4-3-savvas-q24", "Simplify. Watch for difference-of-squares pattern."),
        ("Practice #17 · 4-3-savvas-q17", "Simplify. State all domain restrictions."),
        ("Practice #19 · 4-3-savvas-q19", "Simplify. State all domain restrictions."),
        ("Practice #25 · 4-3-savvas-q25", "Simplify. State all domain restrictions."),
        ("Reinforce · 4-3-savvas-q25", "Extra practice — confirm domain after simplifying."),
    ]
    for lbl, txt in items:
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.65), color=BLUE)
        top += Inches(0.72)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Share / Summary",
           framework_tag="academic conversation", dok=2, minutes=5)
    card(s, "📢 Sentence frame",
         ["\"To simplify a rational expression I first ___, then I cancel ___.",
          " Values excluded from the domain are ___ because ___.\""],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.8), color=BLUE)
    card(s, "🔭 Preview of Period 2",
         ["Next class: multiply and divide rational expressions.",
          "Same factoring moves — now applied across two fractions.",
          "Period 2 ends with a DOK-3 capstone task."],
         Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.3), color=ACCENT)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 1", "Exit Ticket — Summary Recap",
           framework_tag="not graded", dok="1-2", minutes=2)
    card(s, "✍️ Complete on your exit ticket",
         ["1. Today I learned that ___.",
          "2. A factor I can cancel is ___ because ___.",
          "3. One thing I'm still unsure about is ___."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


# ── P2 deck ───────────────────────────────────────────────────────────────

def build_p2(path="L43_P2_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 4-3 · Period 2",
                "Multiply, Divide + ⭐ DOK-3 Closure")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Do Now — Simplify & State Domain",
           framework_tag="bridge from P1", dok=1, minutes=5)
    card(s, "Practice #25 · 4-3-savvas-q25",
         ["Simplify the rational expression from last class.",
          "",
          "State every value excluded from the domain.",
          "Be ready to explain WHY each restriction exists."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Launch — Example 3 (Multiply/Divide)",
           framework_tag="teacher models", dok=2, minutes=12)
    card(s, "📘 Multiply & Divide Rules",
         ["1. Factor everything.",
          "2. Division → multiply by reciprocal.",
          "3. Cancel factors across any num/denom pair.",
          "4. Domain = union of ALL denominator zeros,",
          "   including reciprocal's.",
          "",
          "Example (4-3-savvas-example-3-lesson-4-3): model steps on board.",
          "  Show factoring → flip → cancel → restricted domain."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GREEN)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Explore — TPS + ⭐ DOK-3 Capstone",
           framework_tag="student work", dok="2-3", minutes=33)
    text(s, "8 items in order. Plan ∼15 min for Practice #13 ⭐ DOK-3 Capstone.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Try It 3 · 4-3-savvas-try-it-3-lesson-4-3", "Multiply rational expressions. Factor, cancel, state domain."),
        ("Practice #27 · 4-3-savvas-q27", "Multiply. Factor both pairs before canceling."),
        ("Try It 4 · 4-3-savvas-try-it-4-lesson-4-3", "Divide rational expressions. Flip second fraction first."),
        ("Practice #29 · 4-3-savvas-q29", "Divide. State domain from both original denominators."),
        ("Try It 5 · 4-3-savvas-try-it-5-lesson-4-3", "Mixed multiply/divide expression. Factor everything first."),
        ("Practice #32 · 4-3-savvas-q32", "Multiply/divide chain. Full domain union required."),
        ("Practice #13  ⭐ DOK-3 CAPSTONE · 4-3-savvas-q13",
         "Open-ended: construct and simplify a rational product. Justify domain restrictions."),
        ("Reinforce · 4-3-savvas-q40", "Extension — confirm full domain from a multi-step expression."),
    ]
    for i, (lbl, txt) in enumerate(items):
        color = ACCENT if i == 6 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.78), color=color)
        top += Inches(0.84)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Share / Summary",
           framework_tag="DOK-3 reveal + P3 bridge", dok=2, minutes=5)
    card(s, "📢 Sentence frame · DOK-3 Capstone",
         ["\"My rational expression was ___, which simplifies to ___.",
          " Values excluded from the domain are ___ because ___.\""],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), color=BLUE)
    card(s, "🔭 Bridge to Period 3",
         ["Next class: modeling with rational expressions + assessment-critical simplify.",
          "Period 3 rehearses Topic 4 LEHS Q#2 — same structure as Practice #19.",
          "Know your domain restrictions cold."],
         Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), color=ACCENT)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 2", "Exit Ticket — Summary Recap",
           framework_tag="not graded", dok="1-2", minutes=3)
    card(s, "✍️ Complete on your exit ticket",
         ["1. Today I learned that ___.",
          "2. The step I almost forgot in division was ___.",
          "3. I find the domain by ___."],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


# ── P3 deck ───────────────────────────────────────────────────────────────

def build_p3(path="L43_P3_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, "Lesson 4-3 · Period 3",
                "Modeling + Assessment-Critical Simplify")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Do Now — Simplify & Domain",
           framework_tag="bridge from P2", dok=1, minutes=5)
    card(s, "Practice #11 · 4-3-savvas-q11",
         ["Simplify the rational expression.",
          "",
          "State every excluded value from the domain.",
          "Be ready to explain each restriction."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Launch — Example 2 + Example 6 (paired)",
           framework_tag="teacher models 2 examples", dok=2, minutes=12)
    card(s, "📘 Simplify Reprise · 4-3-savvas-example-2-lesson-4-3",
         ["▶ LEHS Q#4 rehearsal: factor and simplify, state domain.",
          "Teacher models domain-restriction notation step-by-step."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.2), color=GREEN)
    card(s, "📘 Modeling Anchor · 4-3-savvas-example-6-lesson-4-3",
         ["MODELING steps:",
          "  1. Name the ratio in words.",
          "  2. Translate to rational expression with units.",
          "  3. Simplify with restrictions.",
          "  4. Interpret in context.",
          "Teacher models full 4-step process."],
         Inches(0.6), Inches(4.1), Inches(12.1), Inches(3.0), color=GOLD)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Explore — TPS · Assessment Rehearsal",
           framework_tag="student work · pure DOK-2", dok=2, minutes=35)
    text(s, "8 items. Plan ∼8 min for ⭐ Practice #19 (Topic 4 LEHS Q#2 rehearsal).",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Try It 6 · 4-3-savvas-try-it-6-lesson-4-3", "Modeling problem. Apply all 4 steps: name → translate → simplify → interpret."),
        ("Practice #14 · 4-3-savvas-q14", "Simplify. Identify all excluded domain values."),
        ("Practice #22 · 4-3-savvas-q22", "Simplify. Factor numerator and denominator; state restrictions."),
        ("Practice #34 · 4-3-savvas-q34", "Modeling context. Write the rational expression, simplify, interpret."),
        ("Practice #37 · 4-3-savvas-q37", "Simplify and state domain. Watch for multi-factor denominator."),
        ("Practice #19  ⭐ ASSESSMENT-CRITICAL", "4-3-savvas-q19 — Topic 4 LEHS Q#2 rehearsal. Simplify, full domain, justify."),
        ("Reinforce · 4-3-savvas-q39", "Extension — confirm domain restrictions on a modeling expression."),
    ]
    for i, (lbl, txt) in enumerate(items):
        color = ACCENT if i == 5 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.78), color=color)
        top += Inches(0.84)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Share / Summary · Concept Summary",
           framework_tag="LEHS 8-Q preview", dok=2, minutes=5)
    card(s, "📊 Concept Summary — Lesson 4-3",
         ["DOMAIN: Excludes zeros of every denominator in every step,",
          "  including canceled factors.",
          "",
          "MODELING: Name the ratio in words, translate to rationals with units,",
          "  simplify with restrictions, interpret in context.",
          "",
          "\U0001f3af Topic 4 LEHS assessment: expect Q#2 on simplify + domain restrictions."],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.5), color=BLUE)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Period 3", "Exit Ticket — Pre-Assessment Confidence",
           framework_tag="confidence check", dok="1-2", minutes=3)
    card(s, "✍️ Rate yourself 1–5 + name one review item",
         ["• I can simplify a rational expression and state domain restrictions.  Confidence: 1 / 2 / 3 / 4 / 5",
          "• I can translate a real-world ratio into a rational expression.  Confidence: 1 / 2 / 3 / 4 / 5",
          "",
          "One thing I want to review before Topic 4 assessment: ___"],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


if __name__ == "__main__":
    for fn in [build_p1, build_p2, build_p3]:
        p = fn()
        print(f"Built {p}")
