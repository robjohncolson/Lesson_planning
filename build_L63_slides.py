"""L63_P1_Slides.pptx — projection deck for Lesson 6-3 (single-period).

Single deck: Logarithms + DOK-3 spine q15
(Use graph of y=3^x to estimate log_3 50 — exp↔log inverse relationship).

NOTE: PT#58 dropped as spine — not Form-B-aligned. q15 rehearses Form B Q7,
Q10, Q11 through the exp↔log inverse relationship.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import qb

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x0B, 0x3C, 0x6C)
BLUE   = RGBColor(0x1E, 0x6F, 0xC8)
LIGHT  = RGBColor(0xEA, 0xF2, 0xFB)
DARK   = RGBColor(0x11, 0x22, 0x33)
GRAY   = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0xD9, 0x4E, 0x2F)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xC8, 0x8B, 0x14)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def bg(slide, color=LIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def text(slide, body, left, top, width, height, *,
         size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    lines = body.split("\n") if isinstance(body, str) else body
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb


def badge(slide, label, left, top, *, color=NAVY, text_color=WHITE):
    w = Inches(0.3 + 0.11 * len(label))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.35))
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color
    return w


def header(slide, period_label, phase_title, *, framework_tag, dok, minutes):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    text(slide, phase_title, Inches(0.5), Inches(0.18), Inches(10), Inches(0.9),
         size=34, bold=True, color=WHITE)
    text(slide, f"Algebra 2  ·  Lesson 6-3  ·  {period_label}",
         Inches(0.5), Inches(0.82), Inches(9), Inches(0.35),
         size=14, color=RGBColor(0xBB, 0xD0, 0xE6))
    x = Inches(10.8); y = Inches(0.25)
    w = badge(slide, f"DOK {dok}", x, y, color=GOLD)
    x += w + Inches(0.1)
    badge(slide, f"{minutes} min", x, y, color=GREEN)
    x = Inches(10.8); y = Inches(0.72)
    badge(slide, framework_tag, x, y, color=BLUE)


def card(slide, title, body_lines, left, top, width, height, *, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.03
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = color; s.line.width = Pt(1.5)
    text(slide, title, left + Inches(0.25), top + Inches(0.15),
         width - Inches(0.5), Inches(0.5), size=18, bold=True, color=color)
    text(slide, body_lines, left + Inches(0.3), top + Inches(0.75),
         width - Inches(0.6), height - Inches(0.9), size=14, color=DARK)


def title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    text(s, title, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.8),
         size=28, color=RGBColor(0xBB, 0xD0, 0xE6), align=PP_ALIGN.CENTER)
    text(s, "Student-centered · No Blooket",
         Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
         size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


RULES_TEXT = (
    "LOG ↔ EXP:  log_b x = y  ⇔  b^y = x\n"
    "COMMON LOG:  log x = log_10 x\n"
    "NATURAL LOG: ln x  = log_e x\n"
    "KEY:         ln(e^t) = t     (solves continuous growth!)\n"
    "UNDEFINED:   log_b 0 and log_b (negative) — no real answer\n"
    "IDENTITY:    log_b (b^k) = k"
)


# ── P1 deck (single period) ───────────────────────────────────────────────────

def build_p1(path="L63_P1_Slides.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(
        prs,
        "Lesson 6-3 · Quick — Logarithms + ⭐ DOK-3 Spine: q15 (exp↔log inverse via graph)",
        "If I know y = 3^x, how do I estimate log_3 50 — and why does it work?",
    )

    # ── Do Now ────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Do Now — 3 items (log warm-up)",
           framework_tag="bridge", dok="1-2", minutes=7)
    card(s, "6-3-savvas-q14 · Error Analysis: e^t = 98",
         ["A student writes: e^t = 98 → t = 98/e. What error did they make? "
          "What is the correct first step?"],
         Inches(0.6), Inches(1.7), Inches(5.9), Inches(2.0), color=GOLD)
    card(s, "6-3-savvas-q16 · Reasoning: log_4 x < 0",
         ["When is log_4 x negative? Explain using the log ↔ exp definition."],
         Inches(0.6), Inches(3.85), Inches(5.9), Inches(2.0), color=GOLD)
    card(s, "6-3-savvas-q3 · Vocab: common vs. natural log",
         ["Fill in: log x has base ___. ln x has base ___. "
          "They are related because ___.",
          "",
          'Sentence frame: “log x = y means ___, while ln x = y means ___.”'],
         Inches(6.7), Inches(1.7), Inches(6.0), Inches(2.0), color=BLUE)
    card(s, "\U0001f3af SENTENCE FRAME",
         ["log_b x = y means b^___ = x. "
          "So ln(e^t) = ___ because the base of ln is ___."],
         Inches(6.7), Inches(3.85), Inches(6.0), Inches(2.0), color=GREEN)

    # ── Launch ────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Launch — Example 3: Evaluate Logarithms",
           framework_tag="teacher models", dok="1-2", minutes=8)
    card(s, "\U0001f4d8 Log ↔ Exp Rules (mark this on your student packet!)",
         [RULES_TEXT],
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.8), color=GREEN)
    card(s, "Example 3: 6-3-savvas-example-3-lesson-6-3 — Evaluate each logarithm",
         ["log_5 125 = ?  →  5^? = 125  →  answer: 3",
          "log_{1/4} 16 = ?  →  (1/4)^? = 16  →  answer: −2",
          "log_3 0 = ?  →  3^? = 0  →  UNDEFINED (no real answer)",
          "log_2 2^8 = ?  →  identity log_b(b^k) = k  →  answer: 8"],
         Inches(0.6), Inches(4.65), Inches(12.1), Inches(2.5), color=BLUE)

    # ── Explore ───────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Explore — Think-Pair-Share + ⭐ DOK-3 Spine q15",
           framework_tag="student work", dok="2-3", minutes=25)
    text(s, "Start with ⭐ Practice #15 (10 min TPS + 5 min pair presentations). Then #52, #53, #54, #17.",
         Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=16, color=GRAY)
    top = Inches(2.2)
    items = [
        ("Practice #15 · 6-3-savvas-q15  ⭐ DOK-3 SPINE",
         "Higher Order Thinking: Use the graph of y=3^x to estimate the value of log_3 50. "
         "Explain your reasoning. [GRAPH PLACEHOLDER — y=3^x, x∈[-1,4.5], y∈[0,85], gridlines]", True),
        ("Practice #52 · 6-3-savvas-q52",
         "DOK-2 application. Apply log rules.", False),
        ("Practice #53 · 6-3-savvas-q53",
         "DOK-2 application. [45-min variant: skip]", False),
        ("Practice #54 · 6-3-savvas-q54",
         "DOK-2 application. [45-min variant: skip]", False),
        ("Practice #17 · 6-3-savvas-q17  (promoted to core)",
         "Critique: is log_3(1/27) = -3? Use definition to verify or refute. (Supports Form B Q13)", False),
    ]
    for lbl, txt, is_dok3 in items:
        color = ACCENT if is_dok3 else BLUE
        card(s, lbl, [txt], Inches(0.6), top, Inches(12.1), Inches(0.82), color=color)
        top += Inches(0.88)

    # ── Share / Summary ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Share / Summary — DOK-3 Reveal + 6-4 Bridge",
           framework_tag="academic conversation", dok="2", minutes=5)
    card(s, "\U0001f4e2 q15 DOK-3 Spine — Key Steps",
         ["On the graph of y=3^x, locate y=50.",
          "Read off the x-value: between x=3 (3^3=27) and x=4 (3^4=81).",
          "Closer to 3.5. Formally: log_3 50 ≈ 3.56.",
          "",
          "Why does this work?",
          "log_3 50 = x  means  3^x = 50.",
          "The x-value where the exponential graph hits y=50 IS the logarithm.",
          "",
          "Sentence frame: 'I found log_3 50 ≈ ___ by reading the x-value where "
          "the graph hits y=50. This works because ___ and ___ are inverses.'"],
         Inches(0.6), Inches(1.7), Inches(7.5), Inches(5.5), color=BLUE)
    card(s, "\U0001f52d Bridge to Lesson 6-4",
         ["Next: logarithmic FUNCTIONS — graphing y = log_b x.",
          "Domain: x > 0.  Asymptote: x = 0.",
          "Today's log ↔ exp definition is the foundation for graphing."],
         Inches(8.3), Inches(1.7), Inches(4.8), Inches(3.0), color=ACCENT)

    # ── Exit Ticket ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single Period", "Exit Ticket — Summary Recap",
           framework_tag="not graded", dok="1-2", minutes=5)
    card(s, "✍️ Exit Ticket — 4 blanks",
         ["1. log_b x = y means ____.",
          "",
          "2. log_2 32 = ____ because 2^? = 32.",
          "",
          "3. To solve e^t = 6.125, take the ____ of both sides giving t = ____.",
          "",
          "4. One biggest thing I learned today: ________________________"],
         Inches(0.6), Inches(1.8), Inches(12.1), Inches(5.0), color=GOLD)

    prs.save(path)
    return path


if __name__ == "__main__":
    p = build_p1("L63_P1_Slides.pptx")
    print(f"Built {p}")
